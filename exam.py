from dotenv import load_dotenv  # NEW
load_dotenv(override=False)     # NEW – read .env if present, but don't clobber real env
from flask import Flask, render_template, request, send_from_directory
import docx
import fitz  # PyMuPDF
try:
    import pdfplumber
except Exception:
    pdfplumber = None
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
from openai import OpenAI
import hashlib
from flask import g, Response
import unicodedata
import math
import random
from werkzeug.utils import secure_filename
from uuid import uuid4
from threading import Lock
import json
import zipfile
import time
from collections import deque
import logging
import re
import shutil, subprocess, tempfile, os
from concurrent.futures import ThreadPoolExecutor, as_completed
from threading import BoundedSemaphore
_TEX_PAR = int(os.getenv("APP_TEX_PARALLEL", "1"))
_TEX_SEM = BoundedSemaphore(max(1, _TEX_PAR))
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(name)s: %(message)s"
)

# --- Stage 15: progress pruning ---
PROGRESS_TTL_SEC = 24 * 3600
PROGRESS_MAX_ENTRIES = 2000

def _prune_progress():
    now = int(time.time())
    with PROGRESS_LOCK:
        if len(PROGRESS) > PROGRESS_MAX_ENTRIES:
            # drop oldest half
            items = sorted(PROGRESS.items(), key=lambda kv: kv[1].get("ts", 0))
            for k, _ in items[: len(items)//2]:
                PROGRESS.pop(k, None)
        # TTL pass
        dead = [k for k, v in PROGRESS.items() if (now - int(v.get("ts", now))) > PROGRESS_TTL_SEC]
        for k in dead:
            PROGRESS.pop(k, None)

# --- AI LaTeX repair ---

_TECTONIC_CMD = None

# --- Stage 13: security & rate limits (env-togglable) ---
BASIC_AUTH_ENABLED = os.getenv("APP_BASIC_AUTH", "0") in ("1", "true", "yes")
BASIC_AUTH_USER    = os.getenv("APP_USER", "admin")
BASIC_AUTH_PASS    = os.getenv("APP_PASS", "admin")

# Per-IP limits (defaults are conservative; adjust in env if needed)
RATE_UPLOADS_PER_MIN   = int(os.getenv("APP_RATE_UPLOADS_PER_MIN", "6"))   # POST /upload
RATE_STATUS_PER_10S    = int(os.getenv("APP_RATE_STATUS_PER_10S", "50"))   # GET /status
RATE_DOWNLOADS_PER_MIN = int(os.getenv("APP_RATE_DOWNLOADS_PER_MIN", "60"))# GET /download/*

# Internal buckets
_RL = {
    "upload":   {},   # ip -> deque[timestamps]
    "status":   {},
    "download": {},
}
_RL_WINDOW = {
    "upload":   60.0,
    "status":   10.0,
    "download": 60.0,
}
_RL_LOCK = Lock()
UPLOAD_FS_LOCK = Lock()

def env_int(name: str, default: int) -> int:
    try:
        return int(os.getenv(name, str(default)))
    except Exception:
        return default

def env_float(name: str, default: float) -> float:
    try:
        return float(os.getenv(name, str(default)))
    except Exception:
        return default

def env_str(name: str, default: str) -> str:
    return os.getenv(name, default)

# Stage 14: env-overridable limits (defaults unchanged)
MAX_FILES                = env_int("APP_MAX_FILES", 30)
MAX_FILE_MB              = env_int("APP_MAX_FILE_MB", 25)
TOTAL_UPLOAD_MB          = env_int("APP_TOTAL_UPLOAD_MB", 100)

TXT_CHAR_LIMIT           = env_int("APP_TXT_CHAR_LIMIT", 1_000_000)
RTF_CHAR_LIMIT           = env_int("APP_RTF_CHAR_LIMIT", 1_000_000)
DOCX_PARA_LIMIT          = env_int("APP_DOCX_PARA_LIMIT", 50_000)
PPTX_SLIDE_LIMIT         = env_int("APP_PPTX_SLIDE_LIMIT", 2_000)
PDF_PAGE_LIMIT           = env_int("APP_PDF_PAGE_LIMIT", 2_000)
TOTAL_TEXT_CHAR_CAP      = env_int("APP_TOTAL_TEXT_CHAR_CAP", 3_000_000)

ZIP_UNCOMPRESSED_LIMIT_MB= env_int("APP_ZIP_UNCOMP_MB", 300)
ZIP_COMPRESSION_RATIO_MAX= env_float("APP_ZIP_RATIO_MAX", 200.0)

MAX_CONTENT_LENGTH       = TOTAL_UPLOAD_MB * 1024 * 1024  # Flask body cap (kept)

def _detect_tectonic_cmd():
    global _TECTONIC_CMD
    if _TECTONIC_CMD is not None:
        return _TECTONIC_CMD
    if shutil.which("tectonic") is None:
        raise RuntimeError("tectonic not found on PATH.")
    td = tempfile.mkdtemp()
    tex_path = os.path.join(td, "t.tex")
    with open(tex_path, "w", encoding="utf-8") as f:
        f.write("\\documentclass{article}\\begin{document}x\\end{document}")
    try:
        with _TEX_SEM:
            proc = subprocess.run(
                ["tectonic","-X","compile","--outdir",td,tex_path],
                stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=5
            )
        _TECTONIC_CMD = "new" if proc.returncode == 0 else "old"
    except Exception:
        _TECTONIC_CMD = "old"
    finally:
        try:
            shutil.rmtree(td, ignore_errors=True)
        except Exception:
            pass
    return _TECTONIC_CMD


AI_LATEX_SYSTEM_PROMPT = r"""You are a LaTeX normalizer.
Goal: make minimally invasive fixes so the input compiles in LaTeX/MathJax/KaTeX.
Rules:
- Preserve meaning and wording. Do not rewrite sentences.
- Only fix math/markup issues. No rephrasing, no reordering, no deletions.
- Convert stray Unicode math (× ÷ ≤ ≥ α …) to LaTeX macros BUT ONLY inside math.
- Normalize math delimiters: \$...\$ → \(...\), \$\$...\$\$ → \[...\].
- Balance braces and environments; pair \left...\right..; close \begin{..}/\end{..}.
- Keep plain text outside math untouched (including punctuation and quotes).
- Do not introduce new packages or preambles.
- If the input already looks valid, return it unchanged.
- If unsure, prefer leaving text unchanged rather than inventing notation.
Return only the corrected text, no explanations.
"""

_AI_RISK_PATTERNS = [
    r'(?<!\\)\$[^$]*$',
    r'(?<!\\)\$(?:[^$]*\$(?!\s))',
    r'\\left(?!\s*(?:\(|\[|\{|\\langle))',
    r'\\right(?!\s*(?:\)|\]|\}|\\rangle|\.))',
    r'√',
    r'[α-ωΑ-Ωℝℤℕℚℂ∞×÷≤≥≠≈≡∑∏∫∮∈∉∩∪⊂⊆⊃⊇∧∨¬∀∃∴∵]',
    r'(?m)^[^%]*_[A-Za-z0-9](?![^\\]*\\\))',
]
# --- Stage 12: error messages & codes ---
ERR = {
    "invalid_mode": "Invalid mode. Valid option is exam.",
    "missing_title": "Please enter a title.",
    "long_title": "Title too long. Maximum length is 80 characters",
    "missing_qcount": "Please select a number of questions.",
    "too_many_questions": "Number of questions is too large (max is 30).",
    "no_file_part": "No file part",
    "invalid_ext": "Invalid file format: {name}. Please upload .txt, .pdf, .docx, .pptx or .rtf",
    "no_valid_files": "No valid files to process",
    "no_numbered_items": (
        "We couldn't detect any numbered questions in your file(s). "
        "Please ensure items are numbered like '1.' or '2)' etc."
    ),
    "seq_len": "{name} must have exactly {N} items.",
    "invalid_qtype_at": "Invalid question type at position {i}: '{v}'. Valid types are Long, Short, MCQ, Math.",
    "invalid_diff_at": "Invalid difficulty at position {i}: '{v}'. Valid options are easy, medium, hard.",
    "topic_too_long_at": "Additional instructions at position {i} are too long (max {max} characters).",
    "internal": "Internal server error. Please try again.",
}

# Stage 11 additions
ERR.update({
    "too_many_files": f"Too many files uploaded (max is {MAX_FILES}).",
    "file_too_big": f"A file exceeds the per-file size limit ({MAX_FILE_MB} MB).",
    "total_upload_too_big": f"Total upload exceeds {TOTAL_UPLOAD_MB} MB.",
    "mime_mismatch": "File content does not match its extension.",
    "zip_bomb": "Office file appears malformed or overly compressed (possible zip bomb).",
    "pdf_encrypted": "This PDF is password-protected and cannot be processed.",
    "extraction_failed": "Failed to extract text from one or more files.",
})

ERR.update({
    "auth_required": "Authentication required.",
    "auth_invalid":  "Invalid username or password.",
    "rl_upload":     "Rate limit exceeded for uploads. Please wait a moment and try again.",
    "rl_status":     "Too many status checks; please slow down.",
    "rl_download":   "Rate limit exceeded for downloads. Please wait a moment and try again.",
})

_HEADER_PREFIX_RE = re.compile(r"""(?ix)
    ^\s*
    (?:[-–—•*]\s*)?
    (?:q(?:uestion)?\s*\d+\s*[:\-–—.]?\s*)?
    (?:
        long\s*[- ]?\s*answer |
        short\s*[- ]?\s*answer |
        multiple\s*[- ]?\s*choice |
        math (?:\s*/\s*calculation)? |
        calculation
    )
    \s*(?:questions?)?
    (?P<tail>\s*(?:\([^)]{0,60}\)|\[[^\]]{0,60}\])?)
    \s*[:\-–—]*\s*
""")

def chunk_items_by_tokens(items: list[str], max_input_tokens: int) -> list[list[str]]:
    """Greedily pack numbered items into batches without exceeding max_input_tokens."""
    batches, cur, cur_tok = [], [], 0
    for it in items:
        # cost for item content plus numbering/glue
        t = fast_token_estimate(it) + 6
        if cur and (cur_tok + t) > max_input_tokens:
            batches.append(cur); cur, cur_tok = [], 0
        cur.append(it); cur_tok += t
    if cur: batches.append(cur)
    return batches

def strip_category_header_prefix(item: str) -> str:
    m = _HEADER_PREFIX_RE.match(item)
    if not m:
        return item
    tail = m.group('tail') or ''
    rest = item[m.end():]
    prefix = (tail + ' ') if tail and not tail.endswith(' ') else tail
    return (prefix + rest).lstrip()

def strip_headers_from_items(items: list[str]) -> list[str]:
    return [strip_category_header_prefix(it) for it in items]

def normalize_numbering(s: str) -> str:
    s = re.sub(r'(?m)^\s*(\d+)[\)\-:]\s+', r'\1. ', s)
    s = re.sub(r'(?<!\n)\s+(\d+\.\s+)', r'\n\1', s)
    return s

_SPLIT_RE = re.compile(r'(?m)^\s*\d+[\.\)]\s+')

def split_numbered_items(s: str) -> list[str]:
    s = re.sub(r'(?m)^\s*(\d+)[\)\-:]\s+', r'\1. ', s)
    parts = _SPLIT_RE.split(s)
    items = [p.strip() for p in parts if p.strip()]
    return items

def clamp_items(items: list[str], n: int) -> list[str]:
    return items[:n] if len(items) > n else items

RISK_PATTERNS = [re.compile(p) for p in _AI_RISK_PATTERNS]

def ai_fix_latex(text: str, client, model: str = "gpt-4o-mini", max_tokens: int = 2000) -> str:
    try:
        resp = client.chat.completions.create(
            model=model,
            temperature=0,
            max_tokens=max_tokens,
            messages=[
                {"role": "system", "content": AI_LATEX_SYSTEM_PROMPT},
                {"role": "user", "content": text}
            ],
        )
        out = resp.choices[0].message.content or ""
        if 0.5 <= (len(out) / max(1, len(text))) <= 1.5:
            return out
        return text
    except Exception:
        return text

def patch_left_right(text: str) -> str:
    def _balance(pair_open, pair_close, left_cmd=r'\\left', right_cmd=r'\\right'):
        opens = len(re.findall(left_cmd + r'\s*' + re.escape(pair_open), text))
        closes = len(re.findall(right_cmd + r'\s*' + re.escape(pair_close), text))
        t = text
        if opens > closes:
            t += right_cmd + pair_close
        return t
    t = text
    for op, cl in [('(', ')'), ('[', ']'), ('\\{', '\\}')]:
        t = _balance(op, cl)
    return t

LATEX_CHAR_MAP = {
    "×": r"$\times$", "·": r"$\cdot$", "•": r"$\bullet$", "÷": r"$\div$",
    "±": r"$\pm$", "≠": r"$\ne$", "≈": r"$\approx$", "≡": r"$\equiv$",
    "≤": r"$\le$", "≥": r"$\ge$", "<": "<", ">": ">",
    "∞": r"$\infty$", "∝": r"$\propto$",
    "∈": r"$\in$", "∉": r"$\notin$", "∩": r"$\cap$", "∪": r"$\cup$",
    "⊂": r"$\subset$", "⊆": r"$\subseteq$", "⊃": r"$\supset$", "⊇": r"$\supseteq$",
    "∧": r"$\land$", "∨": r"$\lor$", "¬": r"$\lnot$", "∀": r"$\forall$", "∃": r"$\exists$",
    "∴": r"$\therefore$", "∵": r"$\because$",
    "∂": r"$\partial$", "∇": r"$\nabla$", "∮": r"$\oint$",
    "→": r"$\to$", "←": r"$\leftarrow$", "⇒": r"$\Rightarrow$", "⇔": r"$\Leftrightarrow$", "↦": r"$\mapsto$",
    "°": r"$^\circ$", "′": r"$^\prime$", "″": r"$^{\prime\prime}$",
    "…": r"\ldots",
    "ℝ": r"$\mathbb{R}$", "ℤ": r"$\mathbb{Z}$", "ℕ": r"$\mathbb{N}$",
    "ℚ": r"$\mathbb{Q}$", "ℂ": r"$\mathbb{C}$",
    "∑": "", "∏": "", "∫": "",
}

GREEK_MAP = {
    "α": "$\\alpha$", "β": "$\\beta$", "γ": "$\\gamma$", "δ": "$\\delta$", "ε": "$\\varepsilon$",
    "ζ": "$\\zeta$", "η": "$\\eta$", "θ": "$\\theta$", "ι": "$\\iota$", "κ": "$\\kappa$",
    "λ": "$\\lambda$", "μ": "$\\mu$", "ν": "$\\nu$", "ξ": "$\\xi$", "ο": "o",
    "π": "$\\pi$", "ρ": "$\\rho$", "σ": "$\\sigma$", "τ": "$\\tau$", "υ": "$\\upsilon$",
    "φ": "$\\varphi$", "χ": "$\\chi$", "ψ": "$\\psi$", "ω": "$\\omega$",
    "Α": "A", "Β": "B", "Γ": "$\\Gamma$", "Δ": "$\\Delta$", "Ε": "E", "Ζ": "Z", "Η": "H",
    "Θ": "$\\Theta$", "Ι": "I", "Κ": "K", "Λ": "$\\Lambda$", "Μ": "M", "Ν": "N", "Ξ": "$\\Xi$",
    "Ο": "O", "Π": "$\\Pi$", "Ρ": "P", "Σ": "$\\Sigma$", "Τ": "T", "Υ": "$\\Upsilon$",
    "Φ": "$\\Phi$", "Χ": "X", "Ψ": "$\\Psi$", "Ω": "$\\Omega$",
}

FRACTIONS = {
    "½": r"$\tfrac{1}{2}$", "¼": r"$\tfrac{1}{4}$", "¾": r"$\tfrac{3}{4}$",
    "⅓": r"$\tfrac{1}{3}$", "⅔": r"$\tfrac{2}{3}$",
    "⅕": r"$\tfrac{1}{5}$", "⅖": r"$\tfrac{2}{5}$", "⅗": r"$\tfrac{3}{5}$",
    "⅘": r"$\tfrac{4}{5}$", "⅙": r"$\tfrac{1}{6}$", "⅚": r"$\tfrac{5}{6}$",
    "⅛": r"$\tfrac{1}{8}$", "⅜": r"$\tfrac{3}{8}$", "⅝": r"$\tfrac{5}{8}$", "⅞": r"$\tfrac{7}{8}$",
}

SUPERS = dict(zip("⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁼⁽⁾", "0123456789+-=()"))
SUBS = dict(zip("₀₁₂₃₄₅₆₇₈₉₊₋₌₍₎", "0123456789+-=()"))

# ----- math/no-math segmentation -----
_MATH_SEGMENT_PATTERN = re.compile(
    r'(\\\[(?:.|\n)+?\\\])'      # \[...\]
    r'|(\\\((?:.|\n)+?\\\))'     # \(...\)
    r'|(\$\$(?:.|\n)+?\$\$)'     # $$...$$
    r'|(\$(?:.|\n)+?\$)',        # $...$
    re.MULTILINE
)

def _protect_math_segments(text: str):
    placeholders = []
    def repl(m):
        idx = len(placeholders)
        placeholders.append(m.group(0))
        return f"@@MATH{idx}@@"
    protected = _MATH_SEGMENT_PATTERN.sub(repl, text)
    def restore(s):
        for i, seg in enumerate(placeholders):
            s = s.replace(f"@@MATH{i}@@", seg)
        return s
    return protected, restore

def _transform_inside_math(tex: str, fn_disp_and_inl):
    def repl_disp(m):
        inner = m.group(1)
        return r'\[' + fn_disp_and_inl(inner) + r'\]'
    def repl_inl(m):
        inner = m.group(1)
        return r'\(' + fn_disp_and_inl(inner) + r'\)'
    tex = re.sub(r'\\\[(.+?)\\\]', repl_disp, tex, flags=re.S)
    tex = re.sub(r'\\\((.+?)\\\)', repl_inl, tex, flags=re.S)
    return _sanitize_tex_math(tex)

def _transform_inline_math_only(tex: str, fn_inl):
    def repl_inl(m):
        inner = m.group(1)
        return r'\(' + fn_inl(inner) + r'\)'
    return re.sub(r'\\\((.+?)\\\)', repl_inl, tex, flags=re.S)

def _wrap_exponents_outside_math(text: str) -> str:
    protected, restore = _protect_math_segments(text)
    protected = re.sub(r'\b([A-Za-z0-9])\s*\^\s*\(([^)]+)\)', r'$\1^{(\2)}$', protected)
    protected = re.sub(r'\b([A-Za-z0-9])\s*\^\s*([+-]?\d+)\b', r'$\1^{\2}$', protected)
    protected = re.sub(r'\b([A-Za-z0-9])\s*\^\s*([A-Za-z])\b', r'$\1^{\2}$', protected)
    protected = re.sub(r'\b([A-Za-z0-9])\s*\^\s*([+-]?\d+)(?=\s*[A-Za-z])', r'$\1^{\2}$', protected)
    return restore(protected)

def _replace_super_sub_sequences(text: str) -> str:
    def sup_repl(m):
        mapped = "".join(SUPERS.get(ch, "") for ch in m.group(0))
        return f"$^{{{mapped}}}$" if mapped else m.group(0)
    def sub_repl(m):
        mapped = "".join(SUBS.get(ch, "") for ch in m.group(0))
        return f"$_{{{mapped}}}$" if mapped else m.group(0)
    text = re.sub(r"[⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁼⁽⁾]+", sup_repl, text)
    text = re.sub(r"[₀₁₂₃₄₅₆₇₈₉₊₋₌₍₎]+", sub_repl, text)
    return text

def _replace_sqrt(text: str) -> str:
    text = re.sub(r"√\s*\(([^)]+)\)", lambda m: r"$\sqrt{%s}$" % m.group(1), text)
    text = re.sub(r"√\s*([A-Za-z0-9]+)", lambda m: r"$\sqrt{%s}$" % m.group(1), text)
    text = re.sub(r"√(?=\s|$)", lambda m: r"$\sqrt{\quad}$", text)
    return text

def _limits_op(text: str, symbol: str, latex_cmd: str) -> str:
    sym = re.escape(symbol)
    pat1 = rf"{sym}\s*_\s*\{{([^}}]+)\}}\s*\^\s*\{{([^}}]+)\}}"
    text = re.sub(pat1, lambda m: f"${latex_cmd}_{{{m.group(1)}}}^{{{m.group(2)}}}$", text)
    pat2 = rf"{sym}\s*\^\s*\{{([^}}]+)\}}\s*_\s*\{{([^}}]+)\}}"
    text = re.sub(pat2, lambda m: f"${latex_cmd}_{{{m.group(2)}}}^{{{m.group(1)}}}$", text)
    pat3 = rf"{sym}\s*_\s*([A-Za-z0-9+\-*/\\().]+)\s*\^\s*([A-Za-z0-9+\-*/\\().]+)"
    text = re.sub(pat3, lambda m: f"${latex_cmd}_{{{m.group(1)}}}^{{{m.group(2)}}}$", text)
    pat4 = rf"{sym}\s*\^\s*([A-Za-z0-9+\-*/\\().]+)\s*_\s*([A-Za-z0-9+\-*/\\().]+)"
    text = re.sub(pat4, lambda m: f"${latex_cmd}_{{{m.group(2)}}}^{{{m.group(1)}}}$", text)
    pat5 = rf"{sym}([₀₁₂₃₄₅₆₇₈₉]+)([⁰¹²³⁴⁵⁶⁷⁸⁹]*)"
    def subscript_repl(m):
        sub_chars = m.group(1)
        sup_chars = m.group(2) if m.group(2) else ""
        sub_normal = "".join(SUBS.get(ch, ch) for ch in sub_chars)
        sup_normal = "".join(SUPERS.get(ch, ch) for ch in sup_chars)
        if sup_normal:
            return f"${latex_cmd}_{{{sub_normal}}}^{{{sup_normal}}}$"
        else:
            return f"${latex_cmd}_{{{sub_normal}}}$"
    text = re.sub(pat5, subscript_repl, text)
    text = re.sub(rf"\b{re.escape(symbol)}\b", lambda m: f"${latex_cmd}$", text)
    return text

# --- Stage 13: helpers ---

def _client_ip():
    # honor common proxy header; take the first hop
    xff = request.headers.get("X-Forwarded-For", "")
    if xff:
        return xff.split(",")[0].strip()
    return request.remote_addr or "0.0.0.0"

def _auth_ok_for_request():
    """Return True if auth is disabled OR valid Basic creds supplied."""
    if not BASIC_AUTH_ENABLED:
        return True
    auth = request.authorization
    if not auth:
        return False
    return (auth.username == BASIC_AUTH_USER) and (auth.password == BASIC_AUTH_PASS)

def _need_www_auth():
    return (
        "Authentication required",
        401,
        {"WWW-Authenticate": 'Basic realm="ExamApp"'}
    )

def _rate_allow(bucket: str, ip: str, max_count: int) -> bool:
    """Simple sliding-window counter with per-endpoint windows."""
    now = time.time()
    win = _RL_WINDOW[bucket]
    with _RL_LOCK:
        dq = _RL[bucket].setdefault(ip, deque())
        # drop old
        while dq and (now - dq[0]) > win:
            dq.popleft()
        if len(dq) >= max_count:
            return False
        dq.append(now)
        return True

def get_difficulty_profile(difficulty: str):
    """Maps difficulty to summarization + question-gen behavior."""
    d = (difficulty or "medium").lower().strip()
    if d == "easy":
        return {
            "sum_token_scale": 0.80,           # shorter, quicker summaries
            "sum_temp": 0.05,                  # very stable
            "sum_words_each": "12-18",
            "sum_style": (
                "Prefer high-level, generalized bullets; collapse minor exceptions; "
                "avoid heavy technical detail; keep bullets concise and non-redundant."
            ),
            "q_temp": 0.25,                    # simpler/straighter questions
        }
    elif d == "hard":
        return {
            "sum_token_scale": 1.25,           # longer, more detailed summaries
            "sum_temp": 0.15,                  # a touch more variety
            "sum_words_each": "20-35",
            "sum_style": (
                "Prefer specific, technical bullets; include caveats, edge cases, "
                "quantitative values, assumptions, and parameter ranges where relevant."
            ),
            "q_temp": 0.60,                    # harder, more varied questions
        }
    # medium (default)
    return {
        "sum_token_scale": 1.00,
        "sum_temp": 0.10,
        "sum_words_each": "15-25",
        "sum_style": (
            "Balance breadth and depth; include definitions, formulas, representative examples, "
            "and key process steps with essential qualifiers."
        ),
        "q_temp": 0.40,
    }


# ---- NEW: malformed \frac fixers & backslash sanitizers ----
def _fix_malformed_frac_text(tex: str) -> str:
    def fix(inner: str) -> str:
        # \frac{\text{N}{kg}}  -> \frac{\text{N}}{\text{kg}}
        inner = re.sub(
            r'\\frac\{\\text\{([^{}]+)\}\{([^{}]+)\}\}',
            r'\\frac{\\text{\1}}{\\text{\2}}',
            inner
        )
        # \frac{\mathrm{N}{kg}} -> \frac{\mathrm{N}}{\mathrm{kg}}
        inner = re.sub(
            r'\\frac\{\\mathrm\{([^{}]+)\}\{([^{}]+)\}\}',
            r'\\frac{\\mathrm{\1}}{\\mathrm{\2}}',
            inner
        )
        # \frac{N{kg}} -> \frac{N}{kg}
        inner = re.sub(
            r'\\frac\{([A-Za-z0-9\\]+)\{([A-Za-z0-9\\]+)\}\}',
            r'\\frac{\1}{\2}',
            inner
        )
        # \frac{\text{N}}{kg} -> \frac{\text{N}}{\text{kg}}  (units)
        inner = re.sub(
            r'\\frac\{\\text\{([^{}]+)\}\}\{([A-Za-z]+)\}',
            r'\\frac{\\text{\1}}{\\text{\2}}',
            inner
        )
        # \frac{\text{}}{kg} -> \frac{1}{\text{kg}}  (avoid empty numerator causing brace weirdness)
        inner = re.sub(
            r'\\frac\{\\text\{\}\}\{([A-Za-z]+)\}',
            r'\\frac{1}{\\text{\1}}',
            inner
        )

        return inner
    return _transform_inside_math(tex, fix)

def _sanitize_backslashes(tex: str) -> str:
    # remove \\ inside inline math (illegal), keep display math untouched
    tex = _transform_inline_math_only(tex, lambda s: re.sub(r'\\\\+', ' ', s))
    # now outside math: drop trailing \\ and soften other \\ into a space
    protected, restore = _protect_math_segments(tex)
    # trailing at end-of-item/line or before punctuation
    protected = re.sub(r'\s*\\\\\s*(?=(?:$|\n|[)\]\}.,;:]))', ' ', protected)
    # remaining: collapse to single space
    protected = re.sub(r'\\\\+', ' ', protected)
    return restore(protected)

def _context_aware_math(text: str) -> str:
    protected, restore = _protect_math_segments(text)
    protected = re.sub(r'\be\^([A-Za-z0-9\{\}]+)', r'$e^{\1}$', protected)
    return restore(protected)

def convert_slashes_only_inside_math(tex: str) -> str:
    def _fracify(inner: str) -> str:
        return re.sub(
            r'(?<![A-Za-z0-9_/])([A-Za-z0-9\\][A-Za-z0-9\\^_{}]{0,12})\s*/\s*([A-Za-z0-9\\][A-Za-z0-9\\^_{}]{0,12})(?![A-Za-z0-9_/])',
            r'\\frac{\1}{\2}',
            inner
        )
    return _transform_inside_math(tex, _fracify)

def latex_backup_translate(text: str) -> str:
    if not text:
        return text
    text = normalize_for_latex(text)
    text = _wrap_exponents_outside_math(text)
    text = _context_aware_math(text)
    for k, v in FRACTIONS.items(): text = text.replace(k, v)
    for k, v in GREEK_MAP.items(): text = text.replace(k, v)
    for k, v in LATEX_CHAR_MAP.items():
        if v: text = text.replace(k, v)
    text = _limits_op(text, "∬", r"\\iint")
    text = _limits_op(text, "∭", r"\\iiint")
    text = _limits_op(text, "∮", r"\\oint")
    text = _limits_op(text, "∫", r"\\int")
    text = _limits_op(text, "∑", r"\\sum")
    text = _limits_op(text, "∏", r"\\prod")
    text = _replace_sqrt(text)
    text = _replace_super_sub_sequences(text)
    text = re.sub(r"°\s*C", lambda m: r"$^{\circ}$C", text)
    text = re.sub(r'\bdy/dx\b', lambda m: r'$\frac{dy}{dx}$', text)
    text = re.sub(r'\bd/dx\b', lambda m: r'$\frac{d}{dx}$', text)
    # NEW safety: fix malformed \frac and sanitize stray \\ now
    text = _fix_malformed_frac_text(text)
    text = _sanitize_backslashes(text)
    return text

# ---- NEW: math auto-wrapper for naked macros outside math ----
_MATH_MACRO_CORE = r'(?:' + '|'.join([
    r'frac', r'sqrt', r'left', r'right', r'sum', r'prod', r'int', r'iint', r'iiint', r'oint',
    r'cdot', r'times', r'to', r'le', r'ge', r'neq', r'approx', r'equiv',
    r'alpha', r'beta', r'gamma', r'delta', r'varepsilon', r'theta', r'pi', r'phi', r'varphi',
    r'omega', r'infty', r'partial', r'nabla', r'ldots', r'mathbb', r'mathcal', r'mathrm'
]) + r')'

_MATHY_FRAGMENT = re.compile(
    # A single-line fragment (stop before hard terminators) that contains a math macro
    rf'(?P<frag>[^\n$]*\\{_MATH_MACRO_CORE}[^$.\n;:]*)'
)

_DOLLAR_MATH = re.compile(r'\$(.+?)\$', re.S)

# --- Fix vector/hat/bar macros & malformed \frac, INSIDE MATH ONLY ---
_VEC_LIKE = r'(?:vec|hat|bar|tilde|overline|underline|dot|ddot|breve|check|grave|acute)'

def _fix_veclike_args_in_math(tex: str) -> str:
    r"""Ensure \vec x -> \vec{x}, \hat i -> \hat{i}, etc., and nest properly."""
    def fix(inner: str) -> str:
        # \vec x  -> \vec{x}  (and similar for other macros)
        inner = re.sub(rf'\\({_VEC_LIKE})\s*([A-Za-z])\b', r'\\\1{\2}', inner)
        # \vec{\vec x} -> \vec{\vec{x}}  (rare, but makes braces explicit)
        inner = re.sub(rf'\\({_VEC_LIKE})\s*\{{\s*\\({_VEC_LIKE})\s*([A-Za-z])\s*\}}',
                       r'\\\1{\\\2{\3}}', inner)
        return inner
    return _transform_inside_math(tex, fix)

def _fix_frac_forms_in_math(tex: str) -> str:
    """Normalize various broken \frac forms to exactly two arguments."""
    def fix(inner: str) -> str:
        # 1) Whitespace form: \frac a b  -> \frac{a}{b}
        inner = re.sub(r'\\frac\s+([^\s{}]+)\s+([^\s{}]+)', r'\\frac{\1}{\2}', inner)

        # 2) Triple-arg form where first arg ends with a vec-like macro needing an argument:
        #    \frac{d\vec}{r}{dt} -> \frac{d\vec{r}}{dt}
        inner = re.sub(
            rf'\\frac\s*\{{\s*([^{{}}]*\\(?:{_VEC_LIKE}))\s*\}}\s*\{{\s*([^{{}}]+)\s*\}}\s*\{{\s*([^{{}}]+)\s*\}}',
            r'\\frac{\1{\2}}{\3}', inner)

        # 3) Generic triple-arg fallback: \frac{A}{B}{C} -> \frac{A}{B} (drop 3rd to avoid runaway)
        inner = re.sub(
            r'\\frac\s*\{\s*([^{}]+)\s*\}\s*\{\s*([^{}]+)\s*\}\s*\{\s*([^{}]+)\s*\}',
            r'\\frac{\1}{\2}', inner)

        # 4) Very common slip: "\frac{something" with missing closing }... try to close up to end of group
        #    Heuristic: add a } if we see an opening { without a close before a delimiter.
        # (lightweight—your earlier fixers handle most cases)
        return inner
    return _transform_inside_math(tex, fix)
def _fix_text_macros_in_math(tex: str) -> str:
    """
    Ensure \text has a braced argument inside math and repair the very common slips:
      - '\text dm'  -> '\text{dm}'
      - '\frac{\text}{X}' -> '\frac{\text{}}{X}'  (lets later fixers normalize \frac properly)
    """
    def fix(inner: str) -> str:
        # 1) \text dm  -> \text{dm}       (single or two-word tokens)
        inner = re.sub(r'\\text\s+([A-Za-z]+(?:\s+[A-Za-z]+)?)', r'\\text{\1}', inner)

        # 2) \frac{\text}{...}  -> \frac{\text{}}{...}  (balances braces so later frac normalizers can run)
        inner = re.sub(r'\\frac\s*\{\s*\\text\s*\}\s*(\{)', r'\\frac{\\text{}}\\1', inner)

        return inner
    return _transform_inside_math(tex, fix)
# Keep your _wrap_naked_math() from previous message.


def _wrap_naked_math(tex: str) -> str:
    r"""Wraps fragments containing TeX math macros with \( ... \), but leaves existing math alone."""
    protected, restore = _protect_math_segments(tex)

    # 1) Normalize stray $...$ to \( ... \)
    protected = _DOLLAR_MATH.sub(r'\\(\1\\)', protected)

    # 2) Also normalize the common '$\\to$'
    protected = protected.replace('$\\to$', r'\(\to\)')

    # 3) Wrap fragments that contain math macros but are outside math
    def repl(m):
        frag = m.group('frag')
        # If it's already inside \( ... \) (shouldn't be, we protected), skip
        if frag.strip().startswith(r'\(') and frag.strip().endswith(r'\)'):
            return frag
        return r'\(' + frag.strip() + r'\)'

    # Apply repeatedly in case multiple fragments exist on a line
    for _ in range(3):
        protected = _MATHY_FRAGMENT.sub(repl, protected)

    return restore(protected)

# =========================
# Models
# =========================
# Stage 14: allow overriding models via env
summary_model = env_str("OPENAI_MODEL_SUMMARY", "gpt-4o-mini")
main_model    = env_str("OPENAI_MODEL_MAIN",    "chatgpt-4o-latest")
# Dedicated answers model for "question paper" mode (overridable via env)


def adaptive_summary_length(file_tokens: int, target_summary_tokens: int, content_density: float = 1.0) -> int:
    """Adjust summary length based on content density to avoid incomplete summaries"""
    if file_tokens < 1000:  # Very short files
        return min(target_summary_tokens, file_tokens // 2)
    elif file_tokens > 10000:  # Very long files
        return int(target_summary_tokens * 1.2)  # Allow slightly longer summaries
    else:
        return target_summary_tokens


# =========================
# Normalization helpers
# =========================
def unicode_to_ascii(s):
    normalized = unicodedata.normalize("NFKD", s)
    return "".join(c for c in normalized if not unicodedata.combining(c))

def normalize_for_latex(s: str) -> str:
    s = unicodedata.normalize("NFKD", s)
    s = "".join(ch for ch in s if not unicodedata.combining(ch))
    punct_map = {
        "“": "\"", "”": "\"", "„": "\"", "«": "\"", "»": "\"",
        "‘": "'", "’": "'", "‚": "'", "–": "-", "—": "-", "−": "-",
        "…": "...", "\u00A0": " ",
    }
    for k, v in punct_map.items():
        s = s.replace(k, v)
    s = unicode_to_ascii(s)
    return s

# =========================
# Latency models
# =========================
def t_4o_latest_seconds(n_in_tokens: int, n_out_tokens: int) -> float:
    return 0.6 + (n_in_tokens / 40_000) + (n_out_tokens / 80)

def t_4o_mini_summary_seconds(raw_tokens_per_file: int, summary_tokens_per_file: int) -> float:
    return 0.5 + (raw_tokens_per_file / 60_000) + (summary_tokens_per_file / 100)

def t_non_model_seconds(n_files: int, total_raw_tokens: int, total_questions: int) -> float:
    return 3.0 + 1.0 * n_files + 0.00008 * total_raw_tokens + 0.45 * total_questions

# =========================
# Q/A token estimators
# =========================
def estimate_tokens_main_questions(n_long: int, n_short: int, n_mcq: int, n_math: int) -> int:
    BASE, TOK_LONG, TOK_SHORT, TOK_MCQ, TOK_MATH = 20, 180, 80, 120, 160
    return int(BASE + TOK_LONG * n_long + TOK_SHORT * n_short + TOK_MCQ * n_mcq + TOK_MATH * n_math)

def estimate_tokens_main_answers(n_long: int, n_short: int, n_mcq: int, n_math: int) -> int:
    BASE, TOK_LONG, TOK_SHORT, TOK_MCQ, TOK_MATH = 15, 150, 50, 20, 130
    return int(BASE + TOK_LONG * n_long + TOK_SHORT * n_short + TOK_MCQ * n_mcq + TOK_MATH * n_math)

# =========================
# Summarization planning
# =========================
RECOMMENDED_SUMMARY_TOKENS   = env_int("APP_SUMMARY_TOKENS", 350)
SUMMARY_TOKENS_HARD_MIN      = env_int("APP_SUMMARY_MIN", 200)
SUMMARY_TOKENS_HARD_MAX      = env_int("APP_SUMMARY_MAX", 800)

ALWAYS_SAFE_MAIN_Q_INPUT_CAP = env_int("APP_Q_INPUT_CAP", 12_000)
TARGET_MAX_N_OUT_Q           = env_int("APP_Q_OUT_CAP",   4_000)
TARGET_MAX_N_OUT_A           = env_int("APP_A_OUT_CAP",   2_500)
SUM_MIN_K = env_int("APP_SUM_MIN_K", 2)
SUM_MAX_K = env_int("APP_SUM_MAX_K", 4)

TARGET_MID = 42.0  # More aggressive target
LOWER = 40.0
UPPER = 50.0

def estimate_compile_seconds() -> float:
    # Empirical, conservative for two Tectonic runs in parallel
    return 5.0

def _emergency_tex_sanitize(tex: str) -> str:
    # \frac{\sqrt}{A}{B} -> \frac{\sqrt{A}}{B}
    tex = re.sub(
        r'\\frac\s*\{\s*\\sqrt\s*\}\s*\{\s*([^{}]+)\s*\}\s*\{\s*([^{}]+)\s*\}',
        r'\\frac{\\sqrt{\1}}{\2}', tex
    )
    # \frac{\sqrt}{A} -> \sqrt{A}
    tex = re.sub(
        r'\\frac\s*\{\s*\\sqrt\s*\}\s*\{\s*([^{}]+)\s*\}',
        r'\\sqrt{\1}', tex
    )
    # \sqrt x -> \sqrt{x} (inside doc in case any slipped through)
    tex = re.sub(r'\\sqrt\s+([A-Za-z0-9+\-*/().])', r'\\sqrt{\1}', tex)
    return _sanitize_tex_math(tex)

def plan_summarization_sla(
    timings_so_far: dict,
    n_files: int,
    raw_avg_tokens: int,
    n_out_q_cap: int,
    n_out_a_cap: int,
        min_k: int = SUM_MIN_K,  # was 4
        max_k: int = SUM_MAX_K,   # Increased from 6 for better scaling
    hard_min: int = 150,  # Reduced from 200
    hard_max: int = 1200  # Increased from 800 for longer summaries when needed
):
    T_nonmodel = sum(timings_so_far.get(k, 0.0) for k in ("ingest_write","preprocess","token_count"))
    T_QA_pred = t_4o_latest_seconds(n_in_tokens=n_out_q_cap, n_out_tokens=n_out_a_cap)
    T_compile = estimate_compile_seconds()

    T_rem_for_sum = max(4.0, min(25.0, TARGET_MID - T_nonmodel - T_QA_pred - T_compile))

    K_parallel = choose_summary_parallelism(
        n_files, raw_avg_tokens, RECOMMENDED_SUMMARY_TOKENS,
        T_rem_for_sum, min_k=min_k, max_k=max_k
    )
    S_cap = summary_tokens_cap_per_file(
        T_rem_for_sum, n_files, K_parallel, raw_avg_tokens,
        hard_min=hard_min, hard_max=hard_max
    )
    S_use = min(RECOMMENDED_SUMMARY_TOKENS, S_cap)
    return T_rem_for_sum, K_parallel, S_use


def break_even_raw_avg_tokens_per_file(n_files: int, K_parallel: int, S_summary_tokens_per_file: int) -> float:
    if n_files <= 0 or K_parallel <= 0:
        return float('inf')
    C = math.ceil(n_files / K_parallel)
    alpha = (2.0 / 3.0) * (C / n_files)
    numerator = S_summary_tokens_per_file + (40_000 * C / n_files) * (0.5 + S_summary_tokens_per_file / 100.0)
    denominator = max(1e-9, (1.0 - alpha))
    return numerator / denominator

def break_even_raw_avg_tokens_per_file_simple(S_summary_tokens_per_file: int) -> float:
    return 6_000.0 + 121.0 * S_summary_tokens_per_file

def summary_tokens_cap_per_file(T_sum_budget_s: float, n_files: int, K_parallel: int,
                                raw_avg_tokens_per_file: int,
                                hard_min: int = SUMMARY_TOKENS_HARD_MIN,
                                hard_max: int = SUMMARY_TOKENS_HARD_MAX) -> int:
    if n_files <= 0 or K_parallel <= 0:
        return hard_min
    C = math.ceil(n_files / K_parallel)
    s_max = 100.0 * (T_sum_budget_s / C - 0.5 - raw_avg_tokens_per_file / 60_000.0)
    return int(max(hard_min, min(hard_max, math.floor(s_max))))

def choose_summary_parallelism(n_files, raw_avg_tokens_per_file, S_summary_tokens_per_file,
                               T_sum_budget_s, min_k: int = 2, max_k: int = 8) -> int:
    if n_files <= 0 or T_sum_budget_s <= 0:
        return min_k
    wavesize = (0.5 + raw_avg_tokens_per_file / 60_000.0 + S_summary_tokens_per_file / 100.0)
    k = math.ceil(n_files * wavesize / T_sum_budget_s)
    return max(min_k, min(max_k, k))


def enhance_math_content_for_questions(text_for_main: str, num_math: int, total_questions: int) -> str:
    """
    Enhance the input material to better support math question generation
    by emphasizing computational and procedural content over explanatory text.
    """
    if num_math == 0:
        return text_for_main

    math_ratio = num_math / total_questions if total_questions > 0 else 0

    # If this is primarily a math exam, add a brief instruction to focus on computational aspects
    if math_ratio > 0.5:
        math_focus_instruction = """
Note: This exam focuses primarily on mathematical calculations and problem-solving. 
Emphasize: formulas, procedures, computational techniques, numerical methods, and direct problem-solving over lengthy theoretical explanations.
Extract mathematical problems, worked examples, formulas, and calculation methods.

"""
        return math_focus_instruction + text_for_main
    elif math_ratio > 0.3:
        math_focus_instruction = """
Note: Include mathematical calculations and problem-solving questions.
Focus on: formulas, procedures, and computational techniques alongside conceptual understanding.

"""
        return math_focus_instruction + text_for_main

    return text_for_main

def _fix_sqrt_args_in_math(tex: str) -> str:
    """Ensure \\sqrt has a braced argument inside math."""
    def fix(inner: str) -> str:
        # \sqrt x -> \sqrt{x}
        inner = re.sub(r'\\sqrt\s+([A-Za-z0-9+\-*/().])', r'\\sqrt{\1}', inner)
        # Bare \sqrt (no following { or token) -> \sqrt{}
        inner = re.sub(r'\\sqrt(?!\s*\{)', r'\\sqrt{}', inner)
        return inner
    return _transform_inside_math(tex, fix)

def _fix_frac_sqrt_edgecases_in_math(tex: str) -> str:
    r"""Repair \frac{\sqrt}{A}{B} and \frac{\sqrt}{A} edge cases inside math."""
    def fix(inner: str) -> str:
        # \frac{\sqrt}{A}{B} -> \frac{\sqrt{A}}{B}
        inner = re.sub(
            r'\\frac\s*\{\s*\\sqrt\s*\}\s*\{\s*([^{}]+)\s*\}\s*\{\s*([^{}]+)\s*\}',
            r'\\frac{\\sqrt{\1}}{\2}', inner
        )
        # \frac{\sqrt}{A}  ->  \sqrt{A}   (no denominator present; best-effort)
        inner = re.sub(
            r'\\frac\s*\{\s*\\sqrt\s*\}\s*\{\s*([^{}]+)\s*\}',
            r'\\sqrt{\1}', inner
        )
        return inner
    return _transform_inside_math(tex, fix)

# --- Stage 10: downloads metadata -

def _read_run_meta() -> dict | None:
    try:
        with open(META_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return None

def get_quality_answer_instruction():
    return """Create a comprehensive mark scheme for this question paper.

Quality Requirements:
- Provide detailed marking criteria for each question
- Include multiple acceptable answer variations where appropriate
- For long answers: break down into clear marking points with allocated marks
- For short answers: provide complete answers with key terms highlighted
- For MCQ: state correct answer and explain why other options are incorrect
- For math: show complete working with step-by-step solutions

Format:
- Number each answer to match the question paper exactly
- Use clear, concise marking points separated by ';'
- Use inline math (\\( ... \\)); do NOT insert manual line breaks `\\\\`

Focus on accuracy and completeness over brevity.
"""

def fast_token_estimate(text: str) -> int:
    """Faster token estimation using character count heuristic"""
    return len(text) // 3.8  # More accurate ratio for mixed content

def max_input_tokens_for_main_questions(
        n_files: int,
        total_raw_tokens: int,
        total_questions: int,
        n_out_q_tokens: int,
        n_out_a_tokens: int,
        n_summary_calls_parallel: int,
        raw_avg_tokens_per_file: int,
        S_summary_tokens_per_file: int,
        T_total_sla: float = 60.0,
        headroom_s: float = 10.0,
        T_sum_budget_s: float | None = None,
) -> int:
    T_nonmodel = t_non_model_seconds(n_files, total_raw_tokens, total_questions)
    C = math.ceil(max(1, n_files) / max(1, n_summary_calls_parallel))
    per_file_sum_time = t_4o_mini_summary_seconds(raw_avg_tokens_per_file, S_summary_tokens_per_file)
    T_summaries = C * per_file_sum_time if n_files > 0 else 0.0
    if T_sum_budget_s is not None:
        T_summaries = min(T_summaries, T_sum_budget_s)
    T_answers = t_4o_latest_seconds(n_in_tokens=n_out_q_tokens, n_out_tokens=n_out_a_tokens)
    model_budget = T_total_sla - headroom_s - T_nonmodel
    T_left_for_Q = model_budget - T_summaries - T_answers
    N_in_Q_max = 40_000.0 * (T_left_for_Q - 0.6 - (n_out_q_tokens / 80.0))
    return int(max(0, math.floor(N_in_Q_max)))

# --- Stage 9: progress helpers ---
def fail_progress(job: str, *, pct: int = 97, step: int | None = 5,
                  label: str = "An error occurred", http_status: int = 500, msg: str = "Internal error"):
    """
    Mark job as terminally done (so the frontend stops polling) and return (msg, http_status).
    Uses monotonic pct; defaults near the end so the bar smoothly completes.
    """
    try:
        set_progress(job, pct, step=step, label=label, status="done")
    finally:
        return msg, http_status

# --- Stage 7: blueprint-driven mark scheme helpers ---

def _per_item_answer_spec_lines(blueprint: list[dict], start_idx: int = 1, end_idx: int | None = None) -> str:
    """Reference lines to remind the model of each item's Type/Additional instructions/Difficulty."""
    if end_idx is None:
        end_idx = len(blueprint)
    out = []
    for i in range(start_idx, end_idx + 1):
        it = blueprint[i - 1]
        parts = [f"Type: {it.get('type', 'Long')}"]
        if it.get("difficulty"):
            parts.append(f"Difficulty: {it['difficulty']}")
        if it.get("topic"):
            parts.append(f"Additional instructions: {it['topic']}")
        out.append(f"Item {i}: " + "; ".join(parts))
    return "\n".join(out)

def get_quality_answer_instruction_from_blueprint(blueprint: list[dict]) -> str:
    """Wrap your existing answer instruction with a strict 1..N alignment note + per-item refs."""
    N = len(blueprint)
    base = get_quality_answer_instruction()  # reuse your detailed guidance
    return (
        f"{base}\n"
        f"\nStructure:\n"
        f"- Provide answers for items 1..{N} exactly (no extra items, no missing items)\n"
        f"- Number each answer to match the question paper\n"
        f"\nPer-item reference:\n{_per_item_answer_spec_lines(blueprint)}\n"
    )

def continue_mark_scheme_from_blueprint(
    prev_text: str,
    start_idx: int,
    end_idx: int,
    questions_text: str,
    blueprint: list[dict],
) -> str:
    """Continuation prompt that preserves numbering and uses the per-item plan."""
    subplan = _per_item_answer_spec_lines(blueprint, start_idx, end_idx)
    return f"""Continue the mark scheme by writing items {start_idx} to {end_idx} ONLY.

Rules:
- Number answers {start_idx}..{end_idx} to match the question paper
- Use concise marking points separated by ';'
- For MCQ: state the correct option and briefly why others are incorrect
- For math: show essential working; use inline math \\( ... \\); no manual line breaks `\\\\`
- Do not repeat previous items

Per-item reference:
{subplan}

Question paper:
{questions_text}

Previous mark scheme items:
{prev_text}
"""

# =========================
# App + OpenAI
# =========================
UPLOAD_DIR = "uploads"
OUTPUT_DIR = "generated"
PROGRESS = {}
PROGRESS_LOCK = Lock()
CANCELED_JOBS = set()
CANCELED_LOCK = Lock()

def is_canceled(job: str) -> bool:
    with CANCELED_LOCK:
        return job in CANCELED_JOBS

META_PATH = os.path.join(OUTPUT_DIR, "_meta.json")

def _write_run_meta(mode: str, title: str, available: list[str], extra: dict | None = None):
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    meta = {
        "mode": mode,                       # "exam"
        "title": (title or "").strip(),
        "available": sorted(set(available)),# e.g. ["answers"] or ["questions","answers"]
        "timestamp": int(__import__("time").time())
    }
    if extra:
        meta.update(extra)
    try:
        with open(META_PATH, "w", encoding="utf-8") as f:
            json.dump(meta, f, ensure_ascii=False, indent=2)
    except Exception:
        pass  # non-fatal

def set_progress(job: str, pct: int, step: int | None = None, label: str | None = None, status: str = "running"):
    """Monotonic %; safe to call many times."""
    now = int(time.time())
    with PROGRESS_LOCK:
        state = PROGRESS.setdefault(job, {"pct": 0, "step": 0, "label": "", "status": "running", "ts": now})
        state["pct"] = max(int(state.get("pct", 0)), int(pct))
        if step is not None:
            state["step"] = int(step)
        if label is not None:
            state["label"] = str(label)
        state["status"] = status
        state["ts"] = now


# Stage 1: load OpenAI key from env (fail fast if missing)
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "").strip()
if not OPENAI_API_KEY:
    raise RuntimeError("OPENAI_API_KEY not set. Export it in the environment.")
OPENAI_BASE_URL = os.getenv("OPENAI_BASE_URL", "").strip() or None
client = OpenAI(api_key=OPENAI_API_KEY, base_url=OPENAI_BASE_URL) if OPENAI_BASE_URL \
         else OpenAI(api_key=OPENAI_API_KEY)


ALLOWED_EXTENSIONS = {".txt", ".pdf", ".docx", ".pptx", ".rtf"}
VALID_MODES = {"exam"}
DIFF_ALLOWED = {"easy", "medium", "hard"}
# --- Stage 3: qcount → blueprint helpers ---
Q_TYPES = ("Long", "Short", "MCQ", "Math")

def build_default_blueprint(qcount: int) -> list[dict]:
    """
    Round-robin assignment over (Long, Short, MCQ, Math).
    Returns: [{'index': 1, 'type': 'Long'}, ...] length == qcount
    """
    order = ["Long", "Short", "MCQ", "Math"]
    bp = []
    for i in range(max(0, int(qcount))):
        bp.append({"index": i + 1, "type": order[i % len(order)]})
    return bp

def blueprint_from_legacy_counts(num_long: int, num_short: int, num_mcq: int, num_math: int) -> list[dict]:
    """
    Back-compat: materialize a blueprint from legacy counts in the old grouping order
    (all Long then Short then MCQ then Math).
    """
    bp, idx = [], 1
    for t, n in (("Long", num_long), ("Short", num_short), ("MCQ", num_mcq), ("Math", num_math)):
        for _ in range(max(0, int(n))):
            bp.append({"index": idx, "type": t})
            idx += 1
    return bp

def counts_from_blueprint(bp: list[dict]) -> dict:
    counts = {"Long": 0, "Short": 0, "MCQ": 0, "Math": 0}
    for it in bp:
        t = it.get("type")
        if t in counts:
            counts[t] += 1
    return counts

# --- Stage 5: token cap estimation from blueprint ---
def estimate_output_token_caps_from_bp(blueprint: list[dict]) -> tuple[int, int]:
    """
    Returns (n_out_q_cap, n_out_a_cap) using the blueprint's type counts.
    Keeps the same cap policy (>= TARGET_* and *1.4 headroom).
    """
    c = counts_from_blueprint(blueprint)
    n_out_q_est = estimate_tokens_main_questions(c["Long"], c["Short"], c["MCQ"], c["Math"])
    n_out_a_est = estimate_tokens_main_answers(c["Long"], c["Short"], c["MCQ"], c["Math"])
    n_out_q_cap = max(TARGET_MAX_N_OUT_Q, int(n_out_q_est * 1.4))
    n_out_a_cap = max(TARGET_MAX_N_OUT_A, int(n_out_a_est * 1.4))
    return n_out_q_cap, n_out_a_cap

# --- Stage 6: blueprint-driven question prompt helpers ---

def _bp_counts(blueprint: list[dict]) -> dict:
    # light wrapper (readability)
    return counts_from_blueprint(blueprint)

def _difficulty_profile_for_prompt(global_diff: str | None):
    # reuse same difficulty buckets as your earlier helper
    d = (global_diff or "medium").strip().lower()
    if d not in ("easy", "medium", "hard"):
        d = "medium"

    LONG_RANGE  = {"easy": "40-80",  "medium": "50-100", "hard": "70-120"}
    SHORT_RANGE = {"easy": "8-15",   "medium": "10-20",  "hard": "12-25"}

    diff_line_map = {
        "easy":   "Ensure questions are accessible and cover core topics; prioritize clarity over trickiness",
        "medium": "Ensure questions are appropriately challenging and span different topics",
        "hard":   "Ensure questions are quite difficult but also doable, and span different topics",
    }
    diff_guidance_map = {
        "easy": """- Prefer slightly shorter stems; emphasize definitions and core principles
- Reduce multi-step reasoning; avoid contrived scenarios
- Keep calculations and MCQ distractors straightforward""",
        "medium": """- Balance concise conceptual questions with light real-world context
- Include some multi-step reasoning without excessive length
- MCQ distractors should be plausible and conceptually distinct""",
        "hard": """- Prefer longer, scenario-based items with multi-step reasoning
- Emphasize applications, assumptions, caveats, and edge cases
- MCQ distractors should include subtle common misconceptions""",
    }
    return d, LONG_RANGE[d], SHORT_RANGE[d], diff_line_map[d], diff_guidance_map[d]

def _per_item_spec_lines(blueprint: list[dict], start_idx: int = 1, end_idx: int | None = None) -> str:
    """
    Create a concise per-item plan: 'Item i — Type: X; [Difficulty: y]; [Additional instructions: z]'.
    """
    if end_idx is None:
        end_idx = len(blueprint)
    out = []
    for i in range(start_idx, end_idx + 1):
        it = blueprint[i-1]
        parts = [f"Type: {it.get('type','Long')}"]
        if it.get("difficulty"):
            parts.append(f"Difficulty: {it['difficulty']}")
        if it.get("topic"):
            parts.append(f"Additional instructions: {it['topic']}")
        out.append(f"Item {i}: " + "; ".join(parts))
    return "\n".join(out)

def get_quality_question_instruction_from_blueprint(
    blueprint: list[dict],
    global_difficulty: str | None
) -> str:
    """
    Build the full instruction string using the blueprint order and optional per-item
    difficulty/topics. MCQs keep the 'A) ... B) ... C) ... D) ...' format.
    """
    N = len(blueprint)
    counts = _bp_counts(blueprint)
    d, LONG_RANGE, SHORT_RANGE, diff_line, diff_guidance = _difficulty_profile_for_prompt(global_difficulty)

    # math heavy?
    math_heavy = counts["Math"] / max(1, N) > 0.5

    # quality requirements (mirrors your original but itemized)
    if math_heavy:
        quality_requirements = """Quality Requirements:
- Questions should test mathematical understanding and problem-solving skills
- Focus on: calculations, proofs, algebraic manipulation, geometry, and reasoning
- Long-answer may require working; short-answer should be direct numerical/short responses
- MCQ options should include common calculation errors as distractors
- Use inline math \\( ... \\); avoid display math and do not insert manual line breaks `\\\\`"""
    else:
        quality_requirements = """Quality Requirements:
- Questions must test deep understanding, not just recall
- Long-answer should require analysis/evaluation/synthesis
- Short-answer should test key concepts precisely
- MCQ distractors should be plausible and target misconceptions
- Use inline math \\( ... \\); avoid display math and do not insert manual line breaks `\\\\`"""

    global_diff_line = (
        f"Global difficulty: {d} (apply unless a specific item below sets its own Difficulty)."
        if global_difficulty else
        "No global difficulty set; use per-item Difficulty when provided."
    )

    structure_req = f"""Structure Requirements:
- Produce exactly {N} items, numbered '1.' through '{N}.'
- For MCQ items, use the format: 'N. Question... A) option B) option C) option D) option'
- For Long items, target {LONG_RANGE} words; for Short items, target {SHORT_RANGE} words
- Follow the per-item plan below (order and type are mandatory):"""

    per_item_plan = _per_item_spec_lines(blueprint)

    return f"""Create a high-quality mock exam paper from the material below.

{quality_requirements}

{structure_req}
{per_item_plan}

Quality Standards:
- {diff_line}
{diff_guidance}

{global_diff_line}

YOU MUST NOT INCLUDE ANSWERS TO THE QUESTIONS YOU WRITE

Material:
"""

def continue_numbered_list_from_blueprint(
    prev_text: str,
    start_idx: int,
    end_idx: int,
    material: str,
    blueprint: list[dict],
    global_difficulty: str | None,
) -> str:
    """
    Continuation instruction for missing items using the same per-item plan.
    """
    subplan = _per_item_spec_lines(blueprint, start_idx, end_idx)
    d, LONG_RANGE, SHORT_RANGE, diff_line, diff_guidance = _difficulty_profile_for_prompt(global_difficulty)
    return f"""Continue the mock exam by writing items {start_idx} to {end_idx} ONLY.

Rules:
- Number items {start_idx}..{end_idx} (no others)
- For MCQ items, use 'N. Question... A) option B) option C) option D) option'
- Long ~{LONG_RANGE} words; Short ~{SHORT_RANGE} words
- Use inline math \\( ... \\); no display math; no manual line breaks `\\\\`
- Do not repeat earlier items

Per-item plan:
{subplan}

Previous items (for continuity):
{prev_text}

Material:
{material}
"""

# --- Stage 4: per-question overrides helpers ---
MAX_TOPIC_LEN = 200
QTYPE_NORMALIZE = {
    "long": "Long",
    "short": "Short",
    "mcq": "MCQ",
    "math": "Math",
    "calc": "Math",
    "calculation": "Math",
}

def _norm_type(v: str) -> str | None:
    if not v:
        return None
    key = str(v).strip().lower()
    return QTYPE_NORMALIZE.get(key)

def _norm_diff(v: str) -> str | None:
    if not v:
        return None
    key = str(v).strip().lower()
    return key if key in DIFF_ALLOWED else None

def _sanitize_tex_math(src: str) -> str:
    r"""
    Normalize mixed \( … \) and $ … $ math, remove stray $ inside \( … \),
    and fix \frac a b -> \frac{a}{b}.
    """
    # 1) Remove any $ inside \( … \)
    def _strip_dollars_inside(match: re.Match) -> str:
        inner = match.group(1)
        inner = inner.replace('$', '')
        return r"\(" + inner + r"\)"
    src = re.sub(r"\\\((.*?)\\\)", _strip_dollars_inside, src, flags=re.DOTALL)

    # 2) Convert remaining $…$ to \( … \)
    src = re.sub(r"\$(.+?)\$", lambda m: r"\(" + m.group(1) + r"\)", src, flags=re.DOTALL)

    # 3) Fix common \frac forms missing braces
    #    \frac a b  -> \frac{a}{b}
    src = re.sub(r"\\frac\s+([^\s\{\}]+)\s+([^\s\{\}]+)", r"\\frac{\1}{\2}", src)
    #    \frac{a} b -> \frac{a}{b}
    src = re.sub(r"\\frac\{([^{}]+)\}\s+([^\s\{\}]+)", r"\\frac{\1}{\2}", src)
    #    \frac a {b} -> \frac{a}{b}
    src = re.sub(r"\\frac\s+([^\s\{\}]+)\s+\{([^{}]+)\}", r"\\frac{\1}{\2}", src)

    # 4) Drop any remaining solitary $ (unbalanced)
    src = src.replace('$', '')

    return src
# --- END: TeX math sanitizer ---

def _parse_seq_field(raw: str | None) -> list[str] | None:
    """
    Accepts JSON like '["Long","Short"]' OR CSV 'Long,Short'.
    Returns list of strings, stripped. Returns None if raw is falsy.
    """
    if raw is None:
        return None
    s = str(raw).strip()
    if not s:
        return None
    # Try JSON first
    try:
        val = json.loads(s)
        if isinstance(val, list):
            return [str(x).strip() for x in val]
    except Exception:
        pass
    # Fallback: CSV
    return [t.strip() for t in s.split(",")]

def parse_int(x, default=0):
    try:
        return int(str(x).strip())
    except Exception:
        return default
red_flags = ['Multiple Choice', 'Choice', 'Short', 'Long', 'Math/Calculation', 'Calculation', 'Math',
             'MCQ', 'Question', 'Answer', 'Mark Scheme']

def latex_escape(s: str) -> str:
    return (s or "").replace("\\", r"\textbackslash{}") \
        .replace("&", r"\&").replace("%", r"\%").replace("$", r"\$") \
        .replace("#", r"\#").replace("_", r"\_").replace("{", r"\{") \
        .replace("}", r"\}").replace("~", r"\~{}").replace("^", r"\^{}")

PREAMBLE = r"""
\documentclass[12pt]{article}
\usepackage[a4paper,margin=20mm]{geometry}
\usepackage{amsmath,amssymb}
\usepackage{enumitem}
\setlist[enumerate]{itemsep=0.6em, topsep=0.4em}
\newcommand{\paperheader}[1]{\begin{center}\bfseries\LARGE #1\end{center}\vspace{1em}}
\begin{document}
"""

POSTAMBLE = r"\end{document}"

def tex_from_items(items: list[str], title: str) -> str:
    header = r"\paperheader{" + latex_escape(title) + "}\n"
    body = "\\begin{flushleft}\\begin{enumerate}\n" + \
           "".join("\\item " + it + "\n\n" for it in items) + \
           "\\end{enumerate}\\end{flushleft}\n"
    return PREAMBLE + header + body + POSTAMBLE

# --- Stage 8: qpaper helpers (extract items, make a blueprint, render numbered text) ---
TECTONIC_TIMEOUT = env_int("TECTONIC_TIMEOUT", 45)

def compile_tex_with_tectonic(tex_source: str, *, timeout: int | None = None) -> bytes:
    timeout = TECTONIC_TIMEOUT if timeout is None else timeout
    if shutil.which("tectonic") is None:
        raise RuntimeError("tectonic not found on PATH (make sure your venv/bin/Scripts dir is on PATH).")
    with tempfile.TemporaryDirectory() as td:
        tex_path = os.path.join(td, "doc.tex")
        with open(tex_path, "w", encoding="utf-8") as f:
            f.write(tex_source)
        pdf_path = os.path.join(td, "doc.pdf")
        mode = _detect_tectonic_cmd()
        if mode == "new":
            cmd = ["tectonic", "-X", "compile", "--outdir", td, "--keep-logs", tex_path]
        else:
            cmd = ["tectonic", tex_path, "--keep-logs"]
        with _TEX_SEM:
            proc = subprocess.run(cmd, cwd=td, timeout=timeout,
                                  stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        if proc.returncode != 0 or not os.path.exists(pdf_path):
            log_tail = ""
            log_path = os.path.join(td, "doc.log")
            if os.path.exists(log_path):
                with open(log_path, "r", encoding="utf-8", errors="ignore") as lf:
                    log_tail = lf.read()[-4000:]
            raise RuntimeError(
                "Tectonic failed.\n\n"
                f"STDOUT:\n{proc.stdout.decode(errors='ignore')}\n\n"
                f"STDERR:\n{proc.stderr.decode(errors='ignore')}\n\n"
                f"LOG tail:\n{log_tail}\n\n"
                f"Tried:\n{cmd}"
            )
        with open(pdf_path, "rb") as f:
            return f.read()

# near compile_tex_with_tectonic
def compile_tex_with_tectonic_to_path(tex_source: str, out_path: str, *, timeout: int | None = None) -> None:
    timeout = TECTONIC_TIMEOUT if timeout is None else timeout
    if shutil.which("tectonic") is None:
        raise RuntimeError("tectonic not found on PATH.")
    with tempfile.TemporaryDirectory() as td:
        tex_path = os.path.join(td, "doc.tex")
        with open(tex_path, "w", encoding="utf-8") as f:
            f.write(tex_source)
        mode = _detect_tectonic_cmd()
        cmd = (["tectonic","-X","compile","--outdir",td,"--keep-logs",tex_path]
               if mode == "new" else
               ["tectonic", tex_path, "--keep-logs"])
        with _TEX_SEM:
            proc = subprocess.run(cmd, cwd=td, timeout=timeout, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        pdf_src = os.path.join(td, "doc.pdf")
        if proc.returncode != 0 or not os.path.exists(pdf_src):
            log_tail = ""
            log_path = os.path.join(td, "doc.log")
            if os.path.exists(log_path):
                with open(log_path, "r", encoding="utf-8", errors="ignore") as lf:
                    log_tail = lf.read()[-4000:]
            raise RuntimeError(f"Tectonic failed.\n\nSTDOUT:\n{proc.stdout.decode(errors='ignore')}\n\n"
                               f"STDERR:\n{proc.stderr.decode(errors='ignore')}\n\nLOG tail:\n{log_tail}\n\nTried:\n{cmd}")
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        shutil.move(pdf_src, out_path)

def compile_or_repair_to_path(tex_source: str, out_path: str) -> None:
    first = _sanitize_tex_math(tex_source)
    try:
        return compile_tex_with_tectonic_to_path(first, out_path)
    except Exception:
        second = _sanitize_tex_math(first)
        return compile_tex_with_tectonic_to_path(second, out_path)

def compile_or_repair(tex_source: str, *_, **__) -> bytes:
    # First pass: sanitize & compile
    first = _sanitize_tex_math(tex_source)
    try:
        return compile_tex_with_tectonic(first)
    except Exception:
        # Second pass: run sanitizer again (idempotent) in case upstream injected more issues
        second = _sanitize_tex_math(first)
        return compile_tex_with_tectonic(second)


def parallel_map(func, iterable, max_workers=8):
    results = [None] * len(iterable)
    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        futs = {ex.submit(func, i, x): i for i, x in enumerate(iterable)}
        for fut in as_completed(futs):
            idx = futs[fut]
            results[idx] = fut.result()
    return results

# --- Stage 11: content sniffing & zip safety ---

def _looks_pdf(head: bytes) -> bool:
    return head.startswith(b"%PDF-")

def _looks_rtf(head: bytes) -> bool:
    return head.startswith(b"{\\rtf")

def _looks_zip(head: bytes) -> bool:
    return head.startswith(b"PK\x03\x04") or head.startswith(b"PK\x05\x06") or head.startswith(b"PK\x07\x08")

def _office_zip_kind(fp: str) -> str | None:
    """Return 'docx' if it looks like a Word docx, 'pptx' if a PowerPoint; else None."""
    try:
        with zipfile.ZipFile(fp) as z:
            names = set(z.namelist())
            if "[Content_Types].xml" not in names:
                return None
            if any(n.startswith("word/") for n in names):
                return "docx"
            if any(n.startswith("ppt/") for n in names):
                return "pptx"
    except Exception:
        return None
    return None

def _zip_safety_ok(fp: str) -> bool:
    """Basic zip bomb guard: total uncompressed size and ratio check."""
    try:
        with zipfile.ZipFile(fp) as z:
            total_comp = 0
            total_uncomp = 0
            for i in z.infolist():
                total_comp += max(1, i.compress_size)
                total_uncomp += i.file_size
            if total_uncomp > ZIP_UNCOMPRESSED_LIMIT_MB * 1024 * 1024:
                return False
            ratio = float(total_uncomp) / float(total_comp or 1)
            if ratio > ZIP_COMPRESSION_RATIO_MAX:
                return False
            return True
    except Exception:
        return False

def allowed_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    return ext in ALLOWED_EXTENSIONS

def preprocessing(filepath):
    ext = os.path.splitext(filepath)[1].lower()

    def _cap(s: str, limit: int) -> str:
        if not s:
            return ""
        return s[:limit]

    if ext == ".txt":
        # try utf-8, fallback utf-16
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                return _cap(f.read(), TXT_CHAR_LIMIT)
        except UnicodeDecodeError:
            with open(filepath, "r", encoding="utf-16") as f:
                return _cap(f.read(), TXT_CHAR_LIMIT)

    elif ext == ".pdf":
        # prefer PyMuPDF; cap pages
        try:
            parts = []
            doc = fitz.open(filepath)
            n = min(len(doc), PDF_PAGE_LIMIT)
            for i in range(n):
                parts.append(doc.load_page(i).get_text("text") or "")
            return _cap("\n".join(parts), TXT_CHAR_LIMIT)
        except Exception:
            # fallback to pdfplumber if available
            if pdfplumber is None:
                return ""
            out = []
            try:
                with pdfplumber.open(filepath) as pdf:
                    n = min(len(pdf.pages), PDF_PAGE_LIMIT)
                    for i in range(n):
                        out.append(pdf.pages[i].extract_text() or "")
                return _cap("\n".join(out), TXT_CHAR_LIMIT)
            except Exception:
                return ""

    elif ext == ".docx":
        try:
            doc = docx.Document(filepath)
            paras = []
            for i, p in enumerate(doc.paragraphs):
                if i >= DOCX_PARA_LIMIT:
                    break
                paras.append(p.text)
            return _cap("\n".join(paras), TXT_CHAR_LIMIT)
        except Exception:
            return ""

    elif ext == ".pptx":
        try:
            prs = Presentation(filepath)
            out = []
            for i, slide in enumerate(prs.slides):
                if i >= PPTX_SLIDE_LIMIT:
                    break
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        out.append(shape.text)
            return _cap("\n".join(out), TXT_CHAR_LIMIT)
        except Exception:
            return ""

    elif ext == ".rtf":
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                return _cap(rtf_to_text(f.read()), RTF_CHAR_LIMIT)
        except UnicodeDecodeError:
            with open(filepath, "r", encoding="latin-1", errors="ignore") as f:
                return _cap(rtf_to_text(f.read()), RTF_CHAR_LIMIT)
        except Exception:
            return ""

    return ""


def get_response(instruction, file, model=main_model, max_tokens=1024, temperature=0.2, stream_console=False):
    if client is None:
        raise RuntimeError("OpenAI client not initialized. Set OPENAI_API_KEY in your environment.")

    messages = [
        {"role": "system", "content": "You only generate what the instructions say. "
                                      "Use the exact format requested. Do not skip anything. "
                                      "Do not explain yourself. Do not add titles, headings or "
                                      "comments other than the ones requested."},
        {"role": "user", "content": instruction + "\n\n" + (file or "")}
    ]

    if stream_console:
        full_content = ""
        # Streaming version
        with client.chat.completions.stream(
            model=model,
            messages=messages,
            temperature=temperature,
            max_tokens=max_tokens,
            top_p=1.0,
            frequency_penalty=0,
            presence_penalty=0
        ) as stream:
            for event in stream:
                if event.type == "message.delta" and event.delta.content:
                    print(event.delta.content, end="", flush=True)
                    full_content += event.delta.content
            print()  # newline after completion
        return full_content.strip()

    else:
        # Non-streaming version
        response = client.chat.completions.create(
            model=model,
            messages=messages,
            temperature=temperature,
            max_tokens=max_tokens,
            top_p=1.0,
            frequency_penalty=0,
            presence_penalty=0
        )
        return response.choices[0].message.content.strip()

# =========================
# Web app
# =========================
def website():
    app = Flask(__name__)
    app.logger.setLevel(logging.INFO)
    app.config["MAX_CONTENT_LENGTH"] = MAX_CONTENT_LENGTH

    # --- Stage 15: request ID + security headers ---

    @app.before_request
    def _req_ctx():
        # honor inbound X-Request-ID or make one
        rid = request.headers.get("X-Request-ID")
        g.request_id = (rid.strip() if rid else str(uuid4()))

    @app.after_request
    def _secure_headers(resp: Response):
        # propagate request id
        try:
            resp.headers["X-Request-ID"] = g.request_id
        except Exception:
            pass

        # conservative, CSP omitted to avoid breaking inline scripts/styles
        resp.headers.setdefault("X-Content-Type-Options", "nosniff")
        resp.headers.setdefault("X-Frame-Options", "DENY")
        resp.headers.setdefault("Referrer-Policy", "strict-origin-when-cross-origin")
        resp.headers.setdefault("Permissions-Policy", "geolocation=(), microphone=(), camera=()")

        # avoid caching on sensitive routes
        p = request.path or ""
        if p.startswith(("/upload", "/status")):
            resp.headers["Cache-Control"] = "no-store"
            resp.headers["Pragma"] = "no-cache"
            resp.headers["Expires"] = "0"
        return resp

    # --- Stage 13: global guards ---
    # --- Stage 15: health / readiness probes ---

    @app.route("/healthz")
    def healthz():
        # liveness: process is up
        return {"ok": True, "time": int(time.time())}, 200

    @app.route("/readyz")
    def readyz():
        checks = {}
        ok = True

        # tectonic present
        try:
            mode = _detect_tectonic_cmd()
            checks["tectonic"] = {"ok": True, "mode": mode}
        except Exception as e:
            checks["tectonic"] = {"ok": False, "err": str(e)}
            ok = False

        # openai configured (does not call the API)
        checks["openai_key_present"] = bool(os.getenv("OPENAI_API_KEY", "").strip())
        if not checks["openai_key_present"]:
            ok = False

        # output dir writable
        out_ok = True
        try:
            os.makedirs(OUTPUT_DIR, exist_ok=True)
            testp = os.path.join(OUTPUT_DIR, f".w_{uuid4().hex}")
            with open(testp, "w", encoding="utf-8") as f:
                f.write("ok")
            os.remove(testp)
        except Exception as e:
            out_ok = False
            checks["output_dir_err"] = str(e)
            ok = False
        checks["output_dir_writable"] = out_ok

        status = 200 if ok else 503
        return {"ok": ok, "checks": checks, "time": int(time.time())}, status
    @app.before_request
    def _global_security_and_limits():
        # light maintenance
        _prune_progress()
        # Allow static files and favicon without checks
        if request.endpoint in {"static"}:
            return

        path = request.path or ""

        # 1) Basic Auth (skip /upload here so we can complete its spinner using fail_progress)
        if BASIC_AUTH_ENABLED and path not in ("/upload", "/healthz", "/readyz"):
            if not _auth_ok_for_request():
                # plain 401 with WWW-Authenticate
                return _need_www_auth()

        # 2) Lightweight rate limiting for /status and /download
        ip = _client_ip()

        if path.startswith("/status"):
            if not _rate_allow("status", ip, RATE_STATUS_PER_10S):
                # JSON is fine here; frontend ignores occasional errors while polling
                return ERR["rl_status"], 429

        if path.startswith("/download"):
            if not _rate_allow("download", ip, RATE_DOWNLOADS_PER_MIN):
                return ERR["rl_download"], 429
    os.makedirs("templates", exist_ok=True)
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    html_path = os.path.join("templates", "upload.html")
    html_code = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <title>Build a Mock Exam Paper</title>
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <style>
  details#advancedCard.adv-disabled .adv-summary{
  opacity: .55;
  cursor: not-allowed;
  pointer-events: none; /* prevent opening */
}
  /* drag & drop affordances */
.qcard[draggable="true"]{ cursor: grab; }
.qcard.dragging{ opacity:.6; }
.qcard.drop-target{ outline:2px dashed var(--brand); outline-offset:2px; }
/* Advanced question widgets — full-width stacked cards */
.qgrid{
  display: flex;
  flex-direction: column;
  gap: 12px;
}

.qcard{
  width: 100%;
  background: var(--chip);
  border: 1px solid var(--border);
  border-radius: 14px;
  padding: 12px 14px;
  display: flex;
  flex-direction: column;
  transition: border-color .15s, transform .06s ease;
}
.qcard:active{ transform: translateY(1px); }

/* header row inside each card */
.qhead{
  display: flex;
  align-items: center;
  justify-content: space-between;
  gap: 12px;
}

.qleft{ display: flex; align-items: center; gap: 10px; min-width: 0; }
.qright{ opacity: .75; font-size: .9rem; color: var(--muted); } /* placeholder space */

.qbadge{
  width: 42px; height: 42px; border-radius: 12px;
  display: flex; align-items: center; justify-content: center;
  font-weight: 700;
  background: var(--panel);
  border: 1px solid var(--border);
}
.field.is-disabled { opacity: .58; }
.qmeta{ display:flex; flex-direction:column; gap:2px; min-width:0; }
.qtitle{ font-weight:600; white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
.qtype{ font-size:.9rem; color: var(--muted); }

/* reserved space for future per-question options */
.qopts{
  margin-top: 10px;
  min-height: 44px;                  /* leaves vertical room for controls later */
  border: 1px dashed var(--border);  /* subtle placeholder; remove later */
  border-radius: 10px;
  background: color-mix(in oklab, var(--panel) 92%, transparent);
}

/* subtle type tinting on the badge */
.qcard[data-type="Long"]  .qbadge{ background: color-mix(in oklab, var(--brand) 18%, var(--panel)); }
.qcard[data-type="Short"] .qbadge{ background: color-mix(in oklab, var(--ok)    18%, var(--panel)); }
.qcard[data-type="MCQ"]   .qbadge{ background: color-mix(in oklab, var(--warn)  18%, var(--panel)); }
.qcard[data-type="Math"]  .qbadge{ background: color-mix(in oklab, var(--err)   16%, var(--panel)); }
  /* Styled selects */
.select {
  width: 100%;
  -webkit-appearance: none;
  appearance: none;
  background: var(--panel);
  color: var(--text);
  border: 1px solid var(--border);
  border-radius: 10px;
  padding: 10px 36px 10px 12px; /* room for chevron */
  line-height: 1.2;
  cursor: pointer;
  transition: border-color .15s, box-shadow .15s, background-color .15s, transform .06s ease;
  box-shadow: 0 1px 0 rgba(0,0,0,.04);
}
.select:hover { border-color: color-mix(in oklab, var(--border) 60%, var(--brand) 40%); }
.select:active { transform: translateY(1px); }
.select:focus { outline: none; }
.select:focus-visible {
  border-color: var(--brand);
  box-shadow: 0 0 0 3px color-mix(in oklab, var(--brand) 30%, transparent);
}

/* error + disabled states (your JS already toggles .invalid) */
.select.invalid { border-color: var(--err) !important; }
.select:disabled { opacity: .6; cursor: not-allowed; }

/* wrapper draws the chevron */
.select-wrap { position: relative; }
.select-wrap::after{
  content:"";
  position:absolute; right:12px; top:50%;
  width:10px; height:10px; pointer-events:none;
  transform: translateY(-50%) rotate(45deg);
  border-right:2px solid var(--muted);
  border-bottom:2px solid var(--muted);
  transition: transform .15s, border-color .15s;
}
.select-wrap:focus-within::after{
  transform: translateY(-50%) rotate(225deg);
  border-color: var(--brand);
}

/* hide legacy arrows */
select::-ms-expand { display:none; }

/* dark dropdown panels (best-effort; some UAs ignore) */
.select option {
  background: var(--panel);
  color: var(--text);
}
  .radio-inline { display:flex; gap:14px; flex-wrap:wrap; align-items:center; }
.radio-inline label { display:flex; align-items:center; gap:6px; }
  .maker-badge{
  position: fixed;
  right: 14px;
  bottom: 14px;
  z-index: 9999;
  padding: 8px 10px;
  font-size: 0.85rem;
  border-radius: 10px;
  background: var(--panel);
  color: var(--muted);
  border: 1px solid var(--border);
  box-shadow: 0 6px 18px rgba(0,0,0,.12);
  user-select: none;
  -webkit-user-select: none;
}
@media (max-width: 640px){
  .maker-badge{ font-size: 0.8rem; padding: 7px 9px; right: 10px; bottom: 10px; }
}
  .invalid {
  border-color: var(--err) !important;
}
  /* Make disabled-looking buttons that ignore clicks */
.btn[aria-disabled="true"] {
#drop[aria-disabled="true"] { opacity:.55; pointer-events:none; }
  opacity: .55;
  pointer-events: none;
}

.mini-spinner {
  width: 14px; height: 14px; border-radius: 50%;
  border: 2px solid color-mix(in oklab, var(--text) 20%, transparent);
  border-top-color: var(--text);
  display: inline-block; vertical-align: text-bottom;
  margin-right: 8px; animation: spin 1s linear infinite;
}
    textarea{
      width:100%; padding:10px 12px; background:var(--panel); color:var(--text);
      border:1px solid var(--border); border-radius:10px; resize:vertical; min-height:120px;
    }
    :root{
      --bg:#0b0c0f; --panel:#12141a; --muted:#a3a8b3; --text:#e9ecf1; --brand:#4f8cff;
      --ok:#1db954; --warn:#f0b429; --err:#ff5a5f; --border:#1d2230; --chip:#1a1f2b;
    }
    @media (prefers-color-scheme: light){
      :root{
        --bg:#f7f8fb; --panel:#ffffff; --muted:#5f6b7a; --text:#0d1117; --brand:#3b6dde;
        --ok:#17a054; --warn:#c79007; --err:#d5484d; --border:#e4e7ee; --chip:#f1f3f8;
      }
    }
    *{box-sizing:border-box}
    /* Let the document grow and paint the gradient across the whole scroll height */
html{
  min-height:100%;
  background: linear-gradient(180deg, var(--bg), color-mix(in oklab, var(--bg) 80%, #000 20%));
  background-attachment: fixed;
  font-family: ui-sans-serif, system-ui, -apple-system, "Segoe UI", Roboto, Helvetica, Arial,
               "Apple Color Emoji","Segoe UI Emoji";
}

body{
  margin:0;
  min-height:100%;
  background: transparent; /* gradient is on <html> */
  color: var(--text);
}

button, input, select, textarea { font: inherit; }
    .wrap{max-width:980px; margin:32px auto; padding:0 16px;}
    .brand{display:flex; align-items:center; gap:12px; margin-bottom:16px}
    .logo{width:36px;height:36px;border-radius:8px;background:linear-gradient(135deg,var(--brand),#7aa6ff)}
    h1{font-size:1.5rem;margin:0}
    .subtitle{color:var(--muted);margin:4px 0 24px}
    .grid{display:grid; grid-template-columns: 1.2fr .8fr; gap:18px}
    @media (max-width:900px){.grid{grid-template-columns:1fr}}

    .card{
      background:var(--panel); border:1px solid var(--border); border-radius:14px; padding:18px;
      box-shadow: 0 10px 30px rgba(0,0,0,.12);
    }

    /* Progress Section */
    .progress-section {
      margin-top: 20px;
      padding: 18px;
      border-radius: 14px;
      background: var(--panel);
      border: 1px solid var(--border);
      display: none;
    }
    
    .progress-section.show {
      display: block;
    }
    
    .progress-header {
      display: flex;
      align-items: center;
      gap: 12px;
      margin-bottom: 16px;
    }
    
    .progress-bar-container {
      width: 100%;
      height: 8px;
      background: var(--chip);
      border-radius: 4px;
      overflow: hidden;
      margin-bottom: 12px;
    }
    
    .progress-bar {
      height: 100%;
      background: linear-gradient(90deg, var(--brand), #7aa6ff);
      width: 0%;
      transition: width 0.3s ease;
    }
    
    .progress-steps {
      display: grid;
      gap: 8px;
    }
    
    .progress-step {
      display: flex;
      align-items: center;
      gap: 12px;
      padding: 8px 0;
      font-size: 0.9rem;
    }
    
    .step-icon {
      width: 20px;
      height: 20px;
      border-radius: 50%;
      display: flex;
      align-items: center;
      justify-content: center;
      font-size: 0.75rem;
      font-weight: bold;
    }
    
    .step-icon.pending {
      background: var(--chip);
      color: var(--muted);
    }
    
    .step-icon.active {
      background: var(--brand);
      color: white;
    }
    
    .step-icon.complete {
      background: var(--ok);
      color: white;
    }
    
    .step-text {
      flex: 1;
    }
    
    .step-time {
      color: var(--muted);
      font-size: 0.85rem;
    }
    
    .time-estimate {
      background: var(--chip);
      padding: 8px 12px;
      border-radius: 8px;
      text-align: center;
      margin-bottom: 16px;
      font-size: 0.9rem;
    }

    /* Dropzone */
    .drop{
      border:1.5px dashed var(--border); border-radius:14px; padding:24px; text-align:center;
      background:color-mix(in oklab, var(--panel) 90%, var(--chip) 10%); transition: border-color .15s, background .15s;
    }
    .drop.dragover{border-color:var(--brand); background:color-mix(in oklab, var(--chip) 60%, var(--panel) 40%)}
    .drop p{margin:6px 0 0; color:var(--muted)}
    .btn{
      appearance:none; border:1px solid var(--border); background:var(--chip);
      color:var(--text); padding:10px 14px; border-radius:10px; cursor:pointer; font-weight:600;
    }
    .btn.primary{background:var(--brand); border-color:transparent; color:white}
    .btn.danger{ background: var(--err); border-color: transparent; color:#fff; }
    .btn.ghost{background:transparent}
    .btn:disabled{opacity:.6; cursor:not-allowed}
    .btn-row{display:flex; gap:10px; justify-content:center; margin-top:12px}

    /* File list */
    .files{display:grid; gap:10px; margin-top:16px}
    .file{
      display:flex; align-items:center; gap:12px; padding:10px 12px; border:1px solid var(--border);
      border-radius:12px; background:var(--chip);
    }
    .badge{
      font-size:.75rem; padding:4px 8px; border-radius:999px; background:var(--panel); border:1px solid var(--border);
      color:var(--muted);
    }
    .file-main{flex:1; min-width:0}
    .file-name{white-space:nowrap; overflow:hidden; text-overflow:ellipsis}
    .file-meta{font-size:.85rem; color:var(--muted)}
    .file-remove{border:none; background:transparent; color:var(--err); cursor:pointer; font-weight:700}

    /* Form */
    .form-row{display:grid; grid-template-columns: 1fr 1fr; gap:12px}
    .form-row.stacked{grid-template-columns:1fr}
    .field{display:flex; flex-direction:column; gap:6px}
    label{font-weight:600}
    input[type="text"], input[type="number"]{
      width:100%; padding:10px 12px; background:var(--panel); color:var(--text);
      border:1px solid var(--border); border-radius:10px;
    }
    input[type="number"]{appearance: textfield}
    .help{color:var(--muted); font-size:.9rem}

    /* Status */
    .status{margin-top:14px; padding:12px; border-radius:12px; display:none}
    .status.show{display:block}
    .status.ok{background:color-mix(in oklab, var(--ok) 15%, transparent); border:1px solid color-mix(in oklab, var(--ok) 60%, var(--border))}
    .status.err{background:color-mix(in oklab, var(--err) 15%, transparent); border:1px solid color-mix(in oklab, var(--err) 60%, var(--border))}
    .status.warn{background:color-mix(in oklab, var(--warn) 15%, transparent); border:1px solid color-mix(in oklab, var(--warn) 60%, var(--border))}

    /* Footer */
    .footer{margin:26px 0 6px; color:var(--muted); font-size:.9rem; text-align:center}
    .spinner{
      width:18px;height:18px;border-radius:50%; border:3px solid color-mix(in oklab, var(--text) 20%, transparent);
      border-top-color: var(--text); animation:spin 1s linear infinite; display:inline-block; vertical-align:middle; margin-right:8px
    }
    @keyframes spin{to{transform:rotate(360deg)}}
    a.link{color:var(--brand); font-weight:600; text-decoration:none}
    a.link:hover{text-decoration:underline}

    /* Animation for smooth transitions */
    .fade-in {
      animation: fadeIn 0.3s ease-in;
    }
    
    @keyframes fadeIn {
      from { opacity: 0; transform: translateY(10px); }
      to { opacity: 1; transform: translateY(0); }
    }
/* Collapsible Advanced header */
details.card { padding: 0; }                 /* we’ll pad the summary instead */
.adv-summary {
  display: flex; align-items: center; justify-content: space-between;
  gap: 12px; padding: 14px 18px; cursor: pointer; user-select: none;
  list-style: none;
}
.adv-summary::-webkit-details-marker { display: none; }
.adv-summary::marker { content: ""; }

/* chevron */
.adv-summary::after {
  content:""; width:10px; height:10px; transform: rotate(-45deg);
  border-right:2px solid var(--muted); border-bottom:2px solid var(--muted);
  transition: transform .15s, border-color .15s;
}
details[open] .adv-summary::after { transform: rotate(45deg); border-color: var(--brand); }

.adv-title { font-size: 1.1rem; font-weight: 700; }
.adv-count { color: var(--muted); font-size: 0.9rem; }

.adv-body { padding: 0 18px 16px; }
.adv-summary{
  display: flex;
  align-items: center;
  gap: 12px;
  padding: 14px 18px;
  cursor: pointer;
  user-select: none;
  /* remove the old space-between if you had it */
  justify-content: flex-start;   /* <- ensures items don't spread */
}

/* Push the count to the right, just before the chevron */
.adv-count{
  margin-left: auto;             /* <- moves it to the right edge */
  margin-right: 8px;             /* small gap before the chevron */
  color: var(--muted);
  font-size: 0.9rem;
}

/* Keep the chevron pinned at the far right as before */
.adv-summary::after{
  content:"";
  width:10px; height:10px; transform: rotate(-45deg);
  border-right:2px solid var(--muted); border-bottom:2px solid var(--muted);
  transition: transform .15s, border-color .15s;
}
details[open] .adv-summary::after{
  transform: rotate(45deg);
  border-color: var(--brand);
}
/* --- Advanced controls box --- */
.adv-controls{
  margin: 8px 0 12px;
  padding: 12px;
  border: 1px solid var(--border);
  border-radius: 10px;
  background: var(--chip);
}
.adv-checks{
  display: grid;
  gap: 10px;
  margin: 0; padding: 0;
  border: 0;
}
.adv-checks .chk{
  display: flex;
  align-items: center;
  gap: 8px;
}
.sr-only{
  position:absolute; width:1px; height:1px; padding:0; margin:-1px;
  overflow:hidden; clip:rect(0,0,0,0); white-space:nowrap; border:0;
}
/* qopts layout: Topic | Difficulty | Type */
.qopts{
  margin-top: 12px;
  padding: 12px;
  border: 1px dashed var(--border);
  border-radius: 10px;
  background: color-mix(in oklab, var(--panel) 92%, transparent);
}
.qrow{
  display: grid;
  grid-template-columns: 1.2fr 1fr 1.2fr; /* Topic | Diff | Type */
  gap: 12px;
}
@media (max-width: 760px){
  .qrow{ grid-template-columns: 1fr; }
}
.qopt{ display:flex; flex-direction:column; gap:8px; }
.qopt label{ font-weight:600; font-size:.95rem; }
.qopt .muted{ color:var(--muted); font-weight:500; }

.qopt input[type="text"]{
  width:100%; padding:10px 12px; background:var(--panel); color:var(--text);
  border:1px solid var(--border); border-radius:10px;
}

.pill-group{ display:flex; gap:10px; flex-wrap:wrap; }
.pill{
  display:inline-flex; align-items:center; gap:6px; padding:8px 10px;
  border:1px solid var(--border); border-radius:999px; background:var(--chip);
  user-select:none;
}
.pill input{ accent-color: var(--brand); }

/* Visually mute disabled groups */
.qopt.is-disabled{ opacity:.58; }
  </style>
</head>
<body>
  <div class="wrap">
    <div class="brand">
      <div class="logo" aria-hidden="true"></div>
      <div>
        <h1>Build a Mock Exam Paper</h1>
        <div class="subtitle">Upload study files, choose question counts, and generate a polished paper + mark scheme.</div>
      </div>
    </div>

    <div class="grid">
      <!-- Left: Files -->
      <section class="card">
        <h2 style="margin:0 0 8px">Your files</h2>
        <div id="drop" class="drop" tabindex="0" role="button" aria-label="Upload files">
          <strong>Drag &amp; drop</strong> .txt, .pdf, .docx, .pptx, .rtf here
          <p>or</p>
          <div class="btn-row">
            <button class="btn" id="chooseBtn" type="button">Choose files</button>
            <input id="file" type="file" multiple accept=".txt,.pdf,.docx,.pptx,.rtf" style="display:none" />
          </div>
          <p class="help" id="limitHelp">Max {{ max_file_mb|int }} MB per file</p>
        </div>

        <div class="files" id="file-list" aria-live="polite"></div>
        <hr style="border:none; border-top:1px solid var(--border); margin:16px 0">
        
        <div class="field">
          <label for="manualText">Or, if you prefer, describe some of the exam materials yourself:</label>
          <textarea id="manualText" rows="6" placeholder="Paste or type your study material here..."></textarea>
          <div class="help">
            We'll treat this exactly like an uploaded .txt file.
            <span id="manualCount" style="float:right">0 characters</span>
          </div>
        </div>
      </section>

      <!-- Right: Options -->
      <section class="card">
        <h2 style="margin:0 0 8px">Options</h2>
        <div class="field">
  <label for="title">Title</label>
  <input type="text" id="title" maxlength="80" placeholder="e.g. Physics Mock Paper" value="" />
  <div id="titleHelp" class="help">Appears on both PDFs</div>
  <div id="titleHint" class="help" style="margin-top:6px; display:none; color: var(--err);">
    Please enter a title.
  </div>
</div>
<div class="field" id="qcountRow" style="margin-top:10px">
  <label for="qcount">Number of questions</label>
  <div class="select-wrap">
    <select id="qcount" class="select"></select>
  </div>
  <div id="qcountHint" class="help" style="margin-top:6px; display:none; color: var(--err);">
    Please select a number of questions.
  </div>
</div>
        <div id="optDifficulty" class="field" style="margin-top:12px">
  <label for="difficulty">Difficulty level for all questions:</label>
  <div class="radio-inline" role="radiogroup" aria-label="Difficulty level">
    <label><input type="radio" name="difficulty" id="diff-easy" value="easy"> Easy <span class="help">(fastest)</span></label>
    <label><input type="radio" name="difficulty" id="diff-medium" value="medium" checked> Medium <span class="help">(standard)</span></label>
    <label><input type="radio" name="difficulty" id="diff-hard" value="hard"> Hard <span class="help">(slowest)</span></label>
  </div>
  <div class="help">Relative estimates only: Easy ≈ 0.8× time, Medium ≈ 1.0×, Hard ≈ 1.2×.</div>
  <div class="help">Individual question settings can be tweaked in Advanced.</div>
</div>


        <div class="btn-row" style="justify-content:flex-start; margin-top:16px">
          <button class="btn primary" id="submitBtn" type="button" disabled>
            <span id="submitSpinner" class="spinner" style="display:none"></span>
            Generate PDFs
          </button>
          <button class="btn ghost" id="clearBtn" type="button">Clear list</button>
        </div>

        <div id="response" class="status" role="status" aria-live="polite"></div>
      </section>
    </div>

<!-- Advanced -->
<details id="advancedCard" class="card fade-in" style="display:none; margin-top:16px">
  <summary class="adv-summary">
    <span class="adv-title">Advanced</span>
    <span class="adv-count" id="advCount" aria-live="polite"></span>
  </summary>

  <div class="adv-body">
  <!-- NEW: per-question customization toggles -->
  <div class="adv-controls">
    <fieldset class="adv-checks">
      <legend class="sr-only">Per-question customization</legend>

      <label class="chk">
        <input type="checkbox" id="advTopicChk">
        Provide additional instructions per question
      </label>

      <label class="chk">
        <input type="checkbox" id="advDiffChk">
        Customize individual question difficulties
      </label>

      <label class="chk">
        <input type="checkbox" id="advTypeChk">
        Customize individual question types
      </label>
    </fieldset>
  </div>

  <div id="advEmpty" class="help">No questions yet. Pick counts in Options.</div>
  <div id="qGrid" class="qgrid" role="list"></div>
</div>
</details>
<!-- Progress Section -->
    <div id="progressSection" class="progress-section" tabindex="-1" role="region" aria-label="Generation progress">
      <div class="progress-header">
        <h3 style="margin:0; font-size:1.1rem;">Generating your exam paper...</h3>
        <div class="spinner"></div>
      </div>
      
      <div class="time-estimate">
        <strong>Estimated time:</strong> <span id="timeEstimate">30-60 seconds</span>
      </div>
      
      <div class="progress-bar-container">
        <div id="progressBar" class="progress-bar"></div>
      </div>
      
      <div class="progress-steps" id="progressSteps">
        <div class="progress-step">
          <div class="step-icon pending" id="step1">1</div>
          <div class="step-text">Processing and analyzing your files...</div>
          <div class="step-time" id="time1"></div>
        </div>
        <div class="progress-step">
          <div class="step-icon pending" id="step2">2</div>
          <div class="step-text">Creating intelligent summaries...</div>
          <div class="step-time" id="time2"></div>
        </div>
        <div class="progress-step">
          <div class="step-icon pending" id="step3">3</div>
          <div class="step-text">Generating exam questions...</div>
          <div class="step-time" id="time3"></div>
        </div>
        <div class="progress-step">
          <div class="step-icon pending" id="step4">4</div>
          <div class="step-text">Creating marking scheme...</div>
          <div class="step-time" id="time4"></div>
        </div>
        <div class="progress-step">
          <div class="step-icon pending" id="step5">5</div>
          <div class="step-text">Compiling PDFs...</div>
          <div class="step-time" id="time5"></div>
        </div>
      </div>
    </div>
    <!-- Downloads -->
<section id="downloadsCard" class="card fade-in" style="display:none; margin-top:16px">
  <h2 style="margin:0 0 8px">Downloads</h2>
  <div class="btn-row" style="justify-content:flex-start">
    <a id="dlQuestions" class="btn" href="#" aria-disabled="true">
      <span class="mini-spinner" aria-hidden="true"></span>Questions.pdf
    </a>
    <a id="dlAnswers" class="btn" href="#" aria-disabled="true">
      <span class="mini-spinner" aria-hidden="true"></span>Answers.pdf
    </a>
  </div>
  <div class="help">These will activate automatically when ready. You can keep browsing this page.</div>
</section>


    <div class="footer">By uploading, you confirm you have rights to the content. Supported: .txt, .pdf, .docx, .pptx, .rtf</div>
  </div>

  <script>
    const MAX_FILE_MB = {{ max_file_mb|int }};
    const validExts = [".txt",".pdf",".docx",".pptx",".rtf"];

    const drop = document.getElementById('drop');
    const downloadsCard = document.getElementById('downloadsCard');
const dlQuestions = document.getElementById('dlQuestions');
const dlAnswers = document.getElementById('dlAnswers');
    const chooseBtn = document.getElementById('chooseBtn');
    const fileInput = document.getElementById('file');
    const fileList = document.getElementById('file-list');
    const submitBtn = document.getElementById('submitBtn');
    const submitSpinner = document.getElementById('submitSpinner');
    const clearBtn = document.getElementById('clearBtn');
    const responseEl = document.getElementById('response');
    const manualTextEl = document.getElementById('manualText');
    const manualCountEl = document.getElementById('manualCount');
    const progressSection = document.getElementById('progressSection');
const optDifficulty = document.getElementById('optDifficulty');
const titleHelp     = document.getElementById('titleHelp');
const advCountEl    = document.getElementById('advCount');
    const MAX_RETRIES = 0; // Avoid hammering TPM after a 429; user can click again
    const progressBar = document.getElementById('progressBar');
    const titleInput = document.getElementById('title');
const titleHint  = document.getElementById('titleHint');
    const timeEstimate = document.getElementById('timeEstimate');
    let qErrorArmed = false; // show question-count error only after a failed submit
    let titleErrorArmed = false;
    let qcountErrorArmed = false; // show qcount error only after a failed submit
    let serverProgress = 0;      // % from backend
  let statusInterval = null;   // polling timer
  let currentJob = null;       // job id for this run
  let genAbort = null;         // AbortController for /upload
let didCancel = false;       // flag to suppress error UI on cancel
let clearListDefault = null; // original Clear-list handler we can restore
  const SOFT_CAP = 85;  
    let lastProgress = 0; // ensures progress never decreases
    let visibleProgress = 0;         // what we actually render
let genStartTs = 0;              // when generation started
// remember the user's global difficulty to restore later
let savedGlobalDifficulty = null;
const INTRO_HOLD_MS = 1800;      // ignore server jumps for the first 1.8s
const MAX_STEP = 2.2;            // max % the bar may advance per tick
const CATCHUP = 0.35;   
function lockUI(){
  document.body.classList.add('ui-locked');

  // Left column
  chooseBtn.disabled = true;
  chooseBtn.setAttribute('aria-disabled','true');
  fileInput.disabled = true;
  drop.setAttribute('aria-disabled','true');
  manualTextEl.disabled = true;

  // Options
  titleInput.disabled = true;
  qcountInput.disabled = true;
  document.querySelectorAll('input[name="difficulty"]').forEach(r => r.disabled = true);
  optDifficulty?.setAttribute('aria-disabled','true');

  // Advanced: keep the <summary> clickable, but disable all controls inside the panel
  advancedCard
    .querySelectorAll(':where(input, textarea, select, button)')
    .forEach(el => { el.disabled = true; });

  // DO NOT touch the Downloads links; they should remain usable.
}

function unlockUI(){
  document.body.classList.remove('ui-locked');

  // Left column
  chooseBtn.disabled = false;
  chooseBtn.removeAttribute('aria-disabled');
  fileInput.disabled = false;
  drop.removeAttribute('aria-disabled');
  manualTextEl.disabled = false;
  clearBtn.disabled = false;
  clearBtn.removeAttribute('aria-disabled');

  // Options
  titleInput.disabled = false;
  qcountInput.disabled = false;
  document.querySelectorAll('input[name="difficulty"]').forEach(r => r.disabled = false);
  optDifficulty?.removeAttribute('aria-disabled');

  // Advanced controls back on; re-apply your per-question enable/disable rules
  advancedCard
    .querySelectorAll(':where(input, textarea, select, button)')
    .forEach(el => { el.disabled = false; });

  applyAdvToggles();     // respects advTopicChk/advDiffChk/advTypeChk
  refreshSubmitState();  // re-evaluate submit availability
}
function swapClearToCancel(){
  // preserve original handler if not already stored
  if (!clearListDefault) clearListDefault = clearBtn.onclick || clearListDefaultHandler;

  clearBtn.textContent = 'Cancel';
  clearBtn.classList.remove('ghost');
  clearBtn.classList.add('danger');
  clearBtn.removeAttribute('aria-disabled');
  clearBtn.disabled = false;
  clearBtn.onclick = cancelGeneration;
}

function restoreClearAsClear(){
  clearBtn.textContent = 'Clear list';
  clearBtn.classList.remove('danger');
  clearBtn.classList.add('ghost');
  clearBtn.onclick = (clearListDefault || clearListDefaultHandler);
  clearBtn.disabled = false;
  clearBtn.removeAttribute('aria-disabled');
}

async function cancelGeneration(){
  // mark intent and stop the request
  didCancel = true;

  // ask backend to stop current job (best effort)
  try {
    if (currentJob) {
      await fetch('/cancel', {
        method:'POST',
        headers:{ 'Content-Type':'application/json' },
        body: JSON.stringify({ job: currentJob })
      });
    }
  } catch (_) {}

  // abort the in-flight /upload fetch (front-end)
  try { if (genAbort) genAbort.abort(); } catch (_) {}

  // stop progress/status UI
if (statusInterval) { clearInterval(statusInterval); statusInterval = null; }
resetProgress();
deactivateDownloads('Generation canceled. No files were produced.');

  // UI: restore
  restoreClearAsClear();
  submitSpinner.style.display = 'none';
  unlockUI();
  setStatus('warn', 'Generation canceled. Your files and settings are preserved.');

  // remove leave protection + restore title (use globals set below)
  try { if (window._beforeUnloadHandler) window.removeEventListener('beforeunload', window._beforeUnloadHandler); } catch(_) {}
  if (window._originalTitle) document.title = window._originalTitle;
  window._beforeUnloadHandler = null;
  genAbort = null;
}
    let files = [];
    let progressInterval;
    let startTime;
const qcountInput = document.getElementById('qcount');
// --- Advanced widgets model + DnD ---
const advancedCard = document.getElementById('advancedCard');
const qGrid        = document.getElementById('qGrid');   // matches your HTML
qGrid.addEventListener('change', (e) => {
  const t = e.target;
  if (!t.matches('input[type="radio"][name^="type_"]')) return;

  const card = t.closest('.qcard');
  if (!card) return;

  const id = card.dataset.id;
  const item = advModel.find(x => x.id === id);
  if (!item) return;

  item.type = t.value;

  const qtypeEl = card.querySelector('.qtype');
  if (qtypeEl) qtypeEl.textContent = `Type: ${item.type}`;

  card.dataset.type = item.type || '';
  if (!item.type) card.removeAttribute('data-type');
  refreshSubmitState();
});

let advModel = [];
// Persist any in-flight edits from the DOM into advModel (topic, per-question diff)
function captureAdvFormState(){
  qGrid.querySelectorAll('.qcard').forEach(card => {
    const id = card.dataset.id;
    const item = advModel.find(x => x.id === id);
    if (!item) return;

    // topic
    const t = card.querySelector(`#topic_${id}`);
    if (t) item.topic = t.value;

    // per-question difficulty
    const d = card.querySelector(`input[name="diff_${id}"]:checked`);
    if (d) item.diff = d.value;
  });
}
let dragSrcId = null;

const mkId = () =>
  (crypto.randomUUID ? crypto.randomUUID() : 'id_' + Math.random().toString(36).slice(2));
// remove the other two qcountInput.addEventListener('change', ...) lines
qcountInput.addEventListener('change', () => {
  syncAdvToQCount();      // only rebuild when the count actually changes
  refreshSubmitState();   // then (cheap) validation / enablement
});
function fillRange(id, min, max){
  const el = document.getElementById(id);
  el.innerHTML = '';
  for (let i = min; i <= max; i++){
    const opt = document.createElement('option');
    opt.value = String(i);
    opt.textContent = String(i);
    el.appendChild(opt);
  }
  el.value = String(min);
}

// Replace old fillRange usage for qcount with this:
  function populateQCountWithPlaceholder(id, min, max){
    const el = document.getElementById(id);
    el.innerHTML = '';

    // Placeholder option
    const ph = document.createElement('option');
    ph.value = '';
    ph.textContent = 'Select';
    ph.disabled = false;    // keep it clickable so users can open the menu
    ph.selected = true;     // default selection
    el.appendChild(ph);

    for (let i = min; i <= max; i++){
      const opt = document.createElement('option');
      opt.value = String(i);
      opt.textContent = String(i);
      el.appendChild(opt);
    }

    // Ensure the placeholder is the current value
    el.value = '';
  }

  // Populate on load
  populateQCountWithPlaceholder('qcount', 1, 30);
applyAdvToggles(); // already present, but safe to call here too
refreshSubmitState();
// Helper: parse counts safely
    function utf8Bytes(str){
      const enc = new TextEncoder();
      return enc.encode(str);
    }
    
    function manualTextToFile(){
      const text = document.getElementById('manualText').value.trim();
      if (!text) return null;
      const blob = new Blob([text], { type: 'text/plain' });
      return new File([blob], 'manual_input.txt', { type: 'text/plain' });
    }
    
    function bytesToSize(bytes){
      const u = ['B','KB','MB','GB']; let i=0, n=bytes;
      while(n>=1024 && i<u.length-1){ n/=1024; i++; }
      return n.toFixed(n>=10||i===0?0:1)+' '+u[i];
    }
    
    function allowedExtension(filename){
      const idx = filename.lastIndexOf(".");
      if (idx < 0) return false;
      const ext = filename.substring(idx).toLowerCase();
      return validExts.includes(ext);
    }
    
    async function calculateFileHash(file){
      const buf = await file.arrayBuffer();
      const hashBuffer = await crypto.subtle.digest('SHA-256', buf);
      const arr = Array.from(new Uint8Array(hashBuffer));
      return arr.map(b=>b.toString(16).padStart(2,'0')).join('');
    }
    
    function setStatus(type, html){
      responseEl.className = 'status show ' + type;
      responseEl.innerHTML = html;
    }
    
    function clearStatus(){
      responseEl.className = 'status';
      responseEl.innerHTML = '';
    }
function refreshSubmitState(){
  const hasManual = !!manualTextEl.value.trim();
  const titleVal  = titleInput.value.trim();
  const isExam = true;
  const qcountOK  = getQuestionCount() > 0; // true only when not on placeholder

  // Enable submit when there is at least one input (files or manual).
  // We do NOT block the button for qcount here so the user can "attempt"
  // and then we show the red hint.
  const hasAnyInput = (files.length > 0 || hasManual);
  // Block submit unless required per-question radios are filled (when toggled on)
const diffOn = !!document.getElementById('advDiffChk')?.checked;
const typeOn = !!document.getElementById('advTypeChk')?.checked;
const totalQs = getQuestionCount();

// Count checked radios in the advanced grid
const diffsChecked = document.querySelectorAll('.qopt-diff input[type="radio"]:checked').length;
const typesChecked = document.querySelectorAll('.qopt-type input[type="radio"]:checked').length;

// You must fill the radios for whichever checklist(s) are on
let perQuestionOK = true;
if (diffOn) perQuestionOK = perQuestionOK && (diffsChecked === totalQs);
if (typeOn) perQuestionOK = perQuestionOK && (typesChecked === totalQs);

// Final enable/disable
submitBtn.disabled = !(hasAnyInput && perQuestionOK);


  // Title inline error
  if (titleErrorArmed && !titleVal) {
    titleHint.style.display = 'block';
    titleInput.classList.add('invalid');
  } else {
    titleHint.style.display = 'none';
    titleInput.classList.remove('invalid');
    if (titleVal) titleErrorArmed = false;
  }

  // qcount inline hint only AFTER an attempted submit
  const qcountHint = document.getElementById('qcountHint');
  const qcountSel  = document.getElementById('qcount');
  const showQErr   = qcountErrorArmed && isExam && !qcountOK;

  if (qcountHint) qcountHint.style.display = showQErr ? 'block' : 'none';
  if (qcountSel)  qcountSel.classList.toggle('invalid', showQErr);

  // Keep Advanced widget count in sync with qcount
}
    function renderFiles(){
      fileList.innerHTML = '';
      files.forEach(({file, hash})=>{
        const item = document.createElement('div');
        item.className = 'file';
        const ext = file.name.split('.').pop()?.toUpperCase() || '';
        item.innerHTML = `
          <span class="badge" aria-label="File type">${ext}</span>
          <div class="file-main">
            <div class="file-name" title="${file.name}">${file.name}</div>
            <div class="file-meta">${bytesToSize(file.size)}</div>
          </div>
          <button class="file-remove" aria-label="Remove ${file.name}">&times;</button>
        `;
         item.querySelector('.file-remove').onclick = ()=>{
   const idx = files.findIndex(f => f.hash === hash);
  if (idx > -1) {
    files.splice(idx, 1);
    renderFiles();
    refreshSubmitState();
}
        };
        fileList.appendChild(item);
      });
    }
    
    async function acceptFileList(list){
  clearStatus();

  const incoming = Array.isArray(list) ? list : Array.from(list || []);
  const existingHashes = new Set(files.map(f => f.hash));
  const batchHashes = new Set();

  for (const f of incoming){
    if (!allowedExtension(f.name)){
      setStatus('warn', `Unsupported format for <strong>${f.name}</strong>. Allowed: ${validExts.join(', ')}`);
      continue;
    }
    if (f.size > MAX_FILE_MB*1024*1024){
      setStatus('warn', `<strong>${f.name}</strong> is larger than ${MAX_FILE_MB} MB.`);
      continue;
    }

    const hash = await calculateFileHash(f);
    if (existingHashes.has(hash) || batchHashes.has(hash)){
      setStatus('warn', `Duplicate content detected for <strong>${f.name}</strong>. Skipped.`);
      continue;
    }

    files.push({ file: f, hash });
    batchHashes.add(hash);
  }

  renderFiles();
  refreshSubmitState();
}


    function estimateProcessingTime() {
  const isExam = true;
  const totalFiles = files.length + (manualTextEl.value.trim() ? 1 : 0);

  const numQuestions = getQuestionCount();

  let estimatedSeconds = 15;
estimatedSeconds += totalFiles * 3;
estimatedSeconds += numQuestions * 1.5;

  if (estimatedSeconds < 20) return "15-30 seconds";
  if (estimatedSeconds < 40) return "30-45 seconds";
  if (estimatedSeconds < 60) return "45-60 seconds";
  return "60-90 seconds";
}

    function updateProgressStep(stepNumber, status, timeText = '') {
      const stepIcon = document.getElementById(`step${stepNumber}`);
      const stepTime = document.getElementById(`time${stepNumber}`);
      
      stepIcon.className = `step-icon ${status}`;
      if (status === 'complete') {
        stepIcon.innerHTML = '✓';
      } else if (status === 'active') {
        stepIcon.innerHTML = stepNumber;
      }
      
      if (stepTime && timeText) {
        stepTime.textContent = timeText;
      }
    }

    function startProgressSimulation() {
  startTime = Date.now();
  progressSection.classList.add('show', 'fade-in');
    // Auto-scroll to the progress section when generation starts
  // Wait one frame so the browser lays out the newly-shown section, then scroll.
  requestAnimationFrame(() => {
    // a11y: move programmatic focus to the region, without causing a second jump
    if (!progressSection.hasAttribute('tabindex')) progressSection.setAttribute('tabindex', '-1');
    progressSection.focus({ preventScroll: true });

    // smooth scroll the section to the top of the viewport
    progressSection.scrollIntoView({ behavior: 'smooth', block: 'start' });
  });
  timeEstimate.textContent = estimateProcessingTime();

  let currentStep = 1;
  let progress = 0;
  lastProgress = 0; // reset monotonic tracker each run
  genStartTs = Date.now();
visibleProgress = 0;

  const stepTimings = [
    { step: 1, duration: 8000, progress: 15 },
    { step: 2, duration: 15000, progress: 35 },
    { step: 3, duration: 20000, progress: 70 },
    { step: 4, duration: 8000, progress: 85 },
    { step: 5, duration: 7000, progress: 100 }
  ];

  let stepStartTime = Date.now();
  updateProgressStep(1, 'active');

  progressInterval = setInterval(() => {
    const stepElapsed = Date.now() - stepStartTime;

    // Handle step transitions
    const currentStepData = stepTimings[currentStep - 1];
    if (currentStepData && stepElapsed >= currentStepData.duration) {
      const stepDuration = ((Date.now() - stepStartTime) / 1000).toFixed(1);
      updateProgressStep(currentStep, 'complete', `${stepDuration}s`);

      currentStep++;
      if (currentStep <= stepTimings.length) {
        stepStartTime = Date.now();
        updateProgressStep(currentStep, 'active');
      }
    }

    // ===== SMOOTHED HYBRID PROGRESS SECTION =====

// 1) Compute simulated target for the current step
const targetProgress = stepTimings.find(s => s.step === currentStep)?.progress || 0;
const randomFactor = Math.random() * 5 - 2.5; // ±2.5%
let sim = targetProgress + randomFactor;
sim = Math.min(SOFT_CAP, sim); // SOFT_CAP = 85

// 2) For the first INTRO_HOLD_MS, ignore server jumps to avoid the "sudden halfway" effect
const now = Date.now();
const inIntro = (now - genStartTs) < INTRO_HOLD_MS;

// Goal is either calm simulation (during intro) or the higher of sim/server (after intro)
let goal = inIntro ? Math.min(sim, SOFT_CAP) : Math.max(sim, serverProgress);

// Keep the UI below 95% until the server is actually near the end
if (serverProgress < 96) goal = Math.min(95, goal);

// 3) Rate-limit toward the goal (no big jumps)
if (goal > visibleProgress) {
  const gap = goal - visibleProgress;                  // how far we are from where we want to be
  const step = Math.min(MAX_STEP, Math.max(0.6, gap * CATCHUP)); // close a fraction, but cap per tick
  visibleProgress += step;
}

// 4) Monotonic safety and render
visibleProgress = Math.max(visibleProgress, lastProgress);
lastProgress = visibleProgress;
progress = visibleProgress;
progressBar.style.width = `${progress}%`;

// ===== END SMOOTHED HYBRID PROGRESS SECTION =====

  }, 500);
}



    function completeProgress() {
      if (progressInterval) {
        clearInterval(progressInterval);
        progressInterval = null;
      }
      if (statusInterval) { clearInterval(statusInterval); statusInterval = null; }
      
      // Complete all remaining steps
      for (let i = 1; i <= 5; i++) {
        updateProgressStep(i, 'complete');
      }
      
      progressBar.style.width = '100%';
      lastProgress = 100; // <-- add this
      visibleProgress = 100;
      
      // Hide progress section after a brief delay
      setTimeout(() => {
        progressSection.classList.remove('show');
      }, 2000);
    }

    function resetProgress() {
      if (progressInterval) {
        clearInterval(progressInterval);
        progressInterval = null;
      }
      if (statusInterval) { clearInterval(statusInterval); statusInterval = null; }
      
      progressSection.classList.remove('show');
progressBar.style.width = '0%';
lastProgress = 0;
serverProgress = 0;
visibleProgress = 0;    // <-- add
genStartTs = 0;         // <-- add

      
      // Reset all steps
      for (let i = 1; i <= 5; i++) {
        updateProgressStep(i, 'pending');
        const stepIcon = document.getElementById(`step${i}`);
        const stepTime = document.getElementById(`time${i}`);
        stepIcon.innerHTML = i;
        if (stepTime) stepTime.textContent = '';
      }
    }
    function startStatusPolling(job) {
  serverProgress = 0;
  if (statusInterval) { clearInterval(statusInterval); statusInterval = null; }
  statusInterval = setInterval(async () => {
    try {
      const r = await fetch(`/status?job=${encodeURIComponent(job)}`, { cache: 'no-store' });
      if (!r.ok) return;
      const s = await r.json();

      if (typeof s.pct === 'number') {
        serverProgress = Math.max(serverProgress, s.pct);
      }
      // Optional: reflect step labels in UI (kept simple here)
      // if (s.label) { /* you could display s.label somewhere */ }

      if (s.status === 'done' || serverProgress >= 100) {
  // Let the smoothing loop take it to 100; then finalize once we're close
  serverProgress = 100;
  if (visibleProgress >= 99) {
    completeProgress();
    clearInterval(statusInterval);
    statusInterval = null;
  }
}
    } catch (e) {
      // ignore intermittent network errors
    }
  }, 1000);
}

    function showDownloadsPending() {
  downloadsCard.style.display = 'block';

// reset the default helper text each run
const help = downloadsCard.querySelector('.help');
if (help) {
  help.textContent = 'These will activate automatically when ready. You can keep browsing this page.';
}

// disable links + ensure each has a spinner
[dlQuestions, dlAnswers].forEach(a => {
  a.setAttribute('aria-disabled', 'true');
  a.removeAttribute('href');
  a.removeAttribute('target');
  if (!a.querySelector('.mini-spinner')) {
    const s = document.createElement('span');
    s.className = 'mini-spinner';
    s.setAttribute('aria-hidden', 'true');
    a.prepend(s);
  }
});
}

function activateDownloads() {
  // Enable & point directly at your Flask downloads
  dlQuestions.removeAttribute('aria-disabled');
  dlQuestions.href = '/download/questions';
  dlQuestions.target = '_blank';

  dlAnswers.removeAttribute('aria-disabled');
  dlAnswers.href = '/download/answers';
  dlAnswers.target = '_blank';

  // Remove the little spinners now that they’re ready
  const spinners = downloadsCard.querySelectorAll('.mini-spinner');
  spinners.forEach(s => s.remove());
}


    // Drag & drop
    ;['dragenter','dragover'].forEach(evt=>{
      drop.addEventListener(evt, e=>{ e.preventDefault(); e.stopPropagation(); drop.classList.add('dragover'); });
    });
    document.getElementById('advTopicChk')?.addEventListener('change', applyAdvToggles);
document.getElementById('advDiffChk') ?.addEventListener('change', applyAdvToggles);
document.getElementById('advTypeChk') ?.addEventListener('change', applyAdvToggles);

// Also call once on load in case Advanced is already open:
applyAdvToggles();
titleInput.addEventListener('input', refreshSubmitState);
    ;['dragleave','drop'].forEach(evt=>{
      drop.addEventListener(evt, e=>{ e.preventDefault(); e.stopPropagation(); drop.classList.remove('dragover'); });
    });
    drop.addEventListener('drop', e=>{
    if (document.body.classList.contains('ui-locked')) return;
      const dt = e.dataTransfer;
      if (dt?.files?.length) acceptFileList(dt.files);
    });

    chooseBtn.onclick = ()=> fileInput.click();
     fileInput.onchange = async () => {
  const picked = Array.from(fileInput.files); // snapshot, not live
  await acceptFileList(picked);               // wait for hashing, etc.
  fileInput.value = '';                       // now safe to clear
};

    manualTextEl.addEventListener('input', ()=>{
      const len = manualTextEl.value.length;
      manualCountEl.textContent = `${len} character${len===1?'':'s'}`;
      refreshSubmitState();
    });
function deactivateDownloads(reason = '') {
  // disable links and remove any pending spinners
  [dlQuestions, dlAnswers].forEach(a => {
    a.setAttribute('aria-disabled', 'true');
    a.removeAttribute('href');
    a.removeAttribute('target');
  });
  const spinners = downloadsCard.querySelectorAll('.mini-spinner');
  spinners.forEach(s => s.remove());

  if (reason) {
    const help = downloadsCard.querySelector('.help');
    if (help) help.textContent = reason;
  }
}
function clearListDefaultHandler(){
  files = []; renderFiles();
  manualTextEl.value = ''; manualCountEl.textContent = '0 characters';
  qErrorArmed = false;
  advModel = []; renderAdvGrid();    // <-- reset Advanced
  refreshSubmitState(); clearStatus(); resetProgress();
}
clearListDefault = clearListDefaultHandler;   // save original
clearBtn.onclick = clearListDefaultHandler;   // wire it up as default


async function runGenerationOnce() {
  clearStatus();          // no yellow boxes carried over
  resetProgress();        // reset steps, bar, timers

    // Passed validation — now proceed
  // --- NEW: auto-open Advanced on generate ---
    // --- NEW: auto-collapse Advanced on generate ---
  advancedCard.open = false;

  submitBtn.disabled = true;
  submitSpinner.style.display = 'inline-block';
lockUI();
didCancel = false;         // new run
swapClearToCancel();       // turn Clear into red Cancel
  // fresh job id for a truly clean restart
  currentJob = (crypto.randomUUID && crypto.randomUUID()) || Math.random().toString(36).slice(2);

  // show disabled downloads each attempt for consistent UX
  showDownloadsPending();

  // leave-protection + title swap per attempt
  const originalTitle = document.title;
  const beforeUnloadHandler = (e)=>{ e.preventDefault(); e.returnValue=''; };
  window.addEventListener('beforeunload', beforeUnloadHandler);
  document.title = 'Generating… ⏳ ' + (document.getElementById('title').value || 'Mock Exam Paper');
    window._originalTitle = originalTitle;            // expose for cancel
window._beforeUnloadHandler = beforeUnloadHandler;
  // start progress + backend polling
  startProgressSimulation();
  startStatusPolling(currentJob);

  try {
    const fd = new FormData();
    fd.append('job', currentJob);

    // files already in memory
    files.forEach(({file})=> fd.append('file[]', file));

    // optional manual text as a .txt file
    const manualFile = manualTextToFile();
    if (manualFile){
      const buf = await manualFile.arrayBuffer();
      const sha = await crypto.subtle.digest('SHA-256', buf);
      const arr = Array.from(new Uint8Array(sha));
      const manualHash = arr.map(b=>b.toString(16).padStart(2,'0')).join('');
       const currentHashes = new Set(files.map(f => f.hash));
 if (!currentHashes.has(manualHash)) {
   fd.append('file[]', manualFile);
 }
    }

    fd.append('qcount', qcountInput.value || '1');
// Keep legacy fields at 0 so backend doesn’t choke (until server supports qcount)
fd.append('long',  '0');
fd.append('short', '0');
fd.append('mcq',   '0');
fd.append('math',  '0');
    fd.append('title', document.getElementById('title').value || 'Mock Exam Paper');
    fd.append('mode', 'exam');
    let chosenDifficulty = '';
const perQuestionDiff = document.getElementById('advDiffChk')?.checked;

if (!perQuestionDiff) {
  chosenDifficulty = (document.querySelector('input[name="difficulty"]:checked') || {}).value || 'medium';
}
// If per-question is on, we intentionally send an empty string (no global difficulty)
fd.append('difficulty', chosenDifficulty);
// Per-question overrides
const topicsOn = document.getElementById('advTopicChk')?.checked;
const diffOn   = document.getElementById('advDiffChk')?.checked;
const typeOn   = document.getElementById('advTypeChk')?.checked;

if (advModel.length) {
  if (typeOn) {
    const types = advModel.map(it => {
      const sel = document.querySelector(`input[name="type_${it.id}"]:checked`);
      return sel ? sel.value : '';
    });
    fd.append('q_types', JSON.stringify(types));
  }
  if (topicsOn) {
    const topics = advModel.map(it => {
      const inp = document.getElementById(`topic_${it.id}`);
      return inp ? inp.value.trim() : '';
    });
    fd.append('q_topics', JSON.stringify(topics));
  }
  if (diffOn) {
    const diffs = advModel.map(it => {
      const sel = document.querySelector(`input[name="diff_${it.id}"]:checked`);
      return sel ? sel.value : '';
    });
    fd.append('q_diffs', JSON.stringify(diffs));
  }
}

    genAbort = new AbortController();
const res = await fetch('/upload', { method:'POST', body: fd, signal: genAbort.signal });
    const text = await res.text();

    if (!res.ok) {
      // fail this attempt cleanly so caller can restart from the top
      throw new Error(text || 'Generation failed');
    }

    // success path
    completeProgress();
    activateDownloads();
    setStatus('ok', 'Your PDFs are ready. Use the buttons below to download.');

    // cleanup after success
    files = []; renderFiles();
    qErrorArmed = false;
    refreshSubmitState();

    return { ok: true };
  } catch (e) {
  // ensure UI is clean for a full restart (progress reset handled by cancelGeneration if used)
  if (!didCancel) {
  resetProgress();
  clearStatus();
  deactivateDownloads('Generation failed. No files were produced.');
}
  if (e && (e.name === 'AbortError' || didCancel)) {
    return { ok: false, canceled: true };
  }
  return { ok: false, error: e?.message || 'Network error.' };
}
 finally {
 restoreClearAsClear();  // always put the Clear button back
genAbort = null;
    // per-attempt cleanup
    if (statusInterval) { clearInterval(statusInterval); statusInterval = null; }
    window.removeEventListener('beforeunload', beforeUnloadHandler);
    document.title = originalTitle;
    submitSpinner.style.display = 'none';
    unlockUI();
    submitBtn.disabled = (files.length === 0 && !manualTextEl.value.trim());
  }
}
/** Render the Advanced grid from advModel; renumber 1..N every time. */
function getQuestionCount(){
  const v = (qcountInput.value || '').trim();
  if (!v) return 0; // placeholder not chosen
  const n = parseInt(v, 10);
  return Number.isNaN(n) ? 0 : Math.max(1, n);
}
function syncAdvToQCount(){
  const N = getQuestionCount();

  // shrink
  if (advModel.length > N) {
    advModel = advModel.slice(0, N);
  }
  // grow
  while (advModel.length < N) {
  advModel.push({ id: mkId(), topic: '', type: '', diff: '' });
}

  renderAdvGrid();
}
function applyAdvToggles(){
  const topicsOn = document.getElementById('advTopicChk')?.checked;
  const diffOn   = document.getElementById('advDiffChk')?.checked;
  const typeOn   = document.getElementById('advTypeChk')?.checked;

  // Topic inputs
  document.querySelectorAll('.qopt-topic').forEach(box=>{
    box.classList.toggle('is-disabled', !topicsOn);
    const inp = box.querySelector('textarea');
    if (inp) inp.disabled = !topicsOn;
  });

  // Difficulty radios (per-question)
  document.querySelectorAll('.qopt-diff').forEach(box=>{
    box.classList.toggle('is-disabled', !diffOn);
    box.querySelectorAll('input[type="radio"]').forEach(r=> r.disabled = !diffOn);
  });

  // Type radios
  document.querySelectorAll('.qopt-type').forEach(box=>{
    box.classList.toggle('is-disabled', !typeOn);
    box.querySelectorAll('input[type="radio"]').forEach(r=> r.disabled = !typeOn);
  });

  // --- Global difficulty (Options tab) ---
  const optDiffGroup  = document.getElementById('optDifficulty');
  const globalRadios  = optDiffGroup ? optDiffGroup.querySelectorAll('input[name="difficulty"]') : [];

  if (diffOn) {
    // Save current selection ONCE when turning per-question on
    if (savedGlobalDifficulty === null) {
      const sel = Array.from(globalRadios).find(r => r.checked);
      savedGlobalDifficulty = sel ? sel.value : null;
    }
    // Visually blank + lock the group while per-question is active
    globalRadios.forEach(r => { r.checked = false; r.disabled = true; });
    optDiffGroup?.setAttribute('aria-disabled','true');
    optDiffGroup?.classList.add('is-disabled');
  } else {
    // Re-enable the group
    globalRadios.forEach(r => { r.disabled = false; });
    optDiffGroup?.removeAttribute('aria-disabled');
    optDiffGroup?.classList.remove('is-disabled');

    // ✅ Only restore a value if we previously saved one
    if (savedGlobalDifficulty !== null) {
      const target = Array.from(globalRadios).find(r => r.value === savedGlobalDifficulty);
      if (target) {
        Array.from(globalRadios).forEach(r => r.checked = false);
        target.checked = true;
      }
      savedGlobalDifficulty = null; // clear stash after restoring
    }
    // 🔒 No "default to medium" here; leave user’s choice untouched.
  }
  // Re-check submit availability whenever toggles change
refreshSubmitState();
}
function renderAdvGrid() {
  const isExam = true;
  const qcount  = getQuestionCount();
  // Persist current input values before we wipe/rebuild the grid
  captureAdvFormState();

  // Advanced is only usable in exam mode AND when qcount is chosen
  if (!isExam || qcount === 0) {
    advancedCard.classList.add('adv-disabled');
    advancedCard.style.display = isExam ? 'block' : 'none'; // show disabled header in exam mode
    advancedCard.removeAttribute('open');                   // can't open
    const emptyEl = document.getElementById('advEmpty');
    if (emptyEl) emptyEl.style.display = isExam ? 'block' : 'none';
    if (advCountEl) advCountEl.textContent = isExam ? 'Select a number' : '';
    qGrid.innerHTML = '';
    return;
  }

  // From here, Advanced is enabled and may open/close
  advancedCard.classList.remove('adv-disabled');
  advancedCard.style.display = 'block';

  qGrid.innerHTML = '';
  const count = advModel.length;

  // show/hide the whole Advanced card
  const emptyEl = document.getElementById('advEmpty');
  if (!count) {
    if (emptyEl) emptyEl.style.display = 'block';
    if (advCountEl) advCountEl.textContent = '';
    advancedCard.removeAttribute('open'); // keep it closed by default
    return;
  }

  if (emptyEl) emptyEl.style.display = 'none';
  if (advCountEl) advCountEl.textContent = `${count} question${count === 1 ? '' : 's'}`;

  advModel.forEach((item, idx) => {
    const el = document.createElement('div');
    el.className = 'qcard';
    el.dataset.id = item.id;
    el.setAttribute('draggable', 'true');

    el.innerHTML = `
  <div class="qhead">
    <div class="qleft">
      <div class="qbadge" aria-label="Question ${idx+1}">${idx+1}</div>
      <div class="qmeta">
        <div class="qtitle">Question ${idx+1}</div>
        <div class="qtype muted">Type: ${item.type || 'N/A'}</div>
      </div>
    </div>
  </div>

  <div class="qopts" role="group" aria-label="Options for question ${idx+1}">
    <div class="qrow">
      <!-- a) Topic -->
      <div class="qopt qopt-topic is-disabled">
        <label for="topic_${item.id}">Additional instructions:</label>
        <textarea id="topic_${item.id}" placeholder="e.g. Make this a multi-part question about the periodic table" disabled maxlength="200" rows="3"></textarea>
        
      </div>

      <!-- b) Difficulty -->
      <div class="qopt qopt-diff is-disabled">
        <label>Difficulty</label>
        <div class="pill-group" role="radiogroup" aria-label="Difficulty">
          <label class="pill"><input type="radio" name="diff_${item.id}" value="easy"   disabled> Easy</label>
          <label class="pill"><input type="radio" name="diff_${item.id}" value="medium" disabled> Medium</label>
          <label class="pill"><input type="radio" name="diff_${item.id}" value="hard"   disabled> Hard</label>
        </div>
      </div>

      <!-- c) Type -->
      <div class="qopt qopt-type is-disabled">
        <label>Question type</label>
        <div class="pill-group" role="radiogroup" aria-label="Question type">
          <label class="pill"><input type="radio" name="type_${item.id}" value="Long"  ${item.type==='Long'  ? 'checked' : ''} disabled> Long</label>
          <label class="pill"><input type="radio" name="type_${item.id}" value="Short" ${item.type==='Short' ? 'checked' : ''} disabled> Short</label>
          <label class="pill"><input type="radio" name="type_${item.id}" value="MCQ"   ${item.type==='MCQ'   ? 'checked' : ''} disabled> MCQ</label>
          <label class="pill"><input type="radio" name="type_${item.id}" value="Math"  ${item.type==='Math'  ? 'checked' : ''} disabled> Math</label>
        </div>
      </div>
    </div>
  </div>
`;
// Restore Topic value and keep it synced to the model
const topicInput = el.querySelector(`#topic_${item.id}`);
if (topicInput) {
  topicInput.value = item.topic || '';
  topicInput.addEventListener('input', (e) => {
    const m = advModel.find(x => x.id === item.id);
    if (m) m.topic = e.target.value;
  });
}

// Restore per-question Difficulty and keep it synced
if (item.diff) {
  const d = el.querySelector(`input[name="diff_${item.id}"][value="${item.diff}"]`);
  if (d) d.checked = true;
}
el.querySelectorAll(`input[name="diff_${item.id}"]`).forEach(r => {
  r.addEventListener('change', (e) => {
    const m = advModel.find(x => x.id === item.id);
    if (m) m.diff = e.target.value;
    refreshSubmitState();
  });
});
if (item.type) {
  el.dataset.type = item.type;     // enables your .qcard[data-type="…"] badge tint
} else {
  el.removeAttribute('data-type');
}
    el.addEventListener('dragstart', (e) => {
      dragSrcId = item.id;
      el.classList.add('dragging');
      e.dataTransfer.effectAllowed = 'move';
      e.dataTransfer.setData('text/plain', dragSrcId);
    });
    el.addEventListener('dragend', () => {
      dragSrcId = null;
      el.classList.remove('dragging');
      qGrid.querySelectorAll('.drop-target').forEach(n => n.classList.remove('drop-target'));
    });
    el.addEventListener('dragover', (e) => { e.preventDefault(); e.dataTransfer.dropEffect = 'move'; });
    el.addEventListener('dragenter', () => { if (el.dataset.id !== dragSrcId) el.classList.add('drop-target'); });
    el.addEventListener('dragleave', () => { el.classList.remove('drop-target'); });
    el.addEventListener('drop', (e) => {
      e.preventDefault();
      const src = e.dataTransfer.getData('text/plain') || dragSrcId;
const tgt = el.dataset.id;
if (!src || !tgt || src === tgt) return;
const i = advModel.findIndex(x => x.id === src);
const j = advModel.findIndex(x => x.id === tgt);
if (i < 0 || j < 0) return;

// NEW: persist any typed values before we mutate advModel
captureAdvFormState();

[advModel[i], advModel[j]] = [advModel[j], advModel[i]];
renderAdvGrid();
    });

    qGrid.appendChild(el);
  });

  // Enable/disable per the three checkboxes
  applyAdvToggles();
  refreshSubmitState();
}
    async function submitFiles(){
  clearStatus();
  resetProgress();

  const isExam = true;
  const titleVal = titleInput.value.trim();
  const qcountOK = getQuestionCount() > 0;

  // Validate BEFORE disabling button / showing spinner
  if (!titleVal) {
    titleErrorArmed = true;
    refreshSubmitState();
    return;
  }

  if (isExam && !qcountOK) {
    qcountErrorArmed = true;       // arm the red complaint now
    refreshSubmitState();          // shows the red hint + invalid border
    document.getElementById('qcount').focus();
    return;
  }

  // Passed validation — now proceed
  submitBtn.disabled = true;
  submitSpinner.style.display = 'inline-block';
// Hard-stop if total questions is 0 (extra safety)

// --- auto-retry wrapper ---
let attempt = 0;
let result;

do {
  attempt++;
  result = await runGenerationOnce();
  if (result.ok) break; // success, stop retrying
} while (attempt <= MAX_RETRIES);

// If still failing after retries, show one final error box (not the inline field hints)
if (!result.ok) {
  if (result.canceled) {
    // user canceled: no error box
  } else {
    setStatus(
      'err',
      `Sorry, we couldn't complete your generation${
        MAX_RETRIES ? ` after ${attempt} attempt(s)` : ''
      }. Please try again later.`
    );
  }
}
    }
    
    submitBtn.onclick = submitFiles;

    // Keyboard activation for drop area
    drop.addEventListener('keydown', e=>{
      if (e.key === 'Enter' || e.key === ' ') { e.preventDefault(); fileInput.click(); }
    });
    renderAdvGrid();
refreshSubmitState();
  </script>
  <div class="maker-badge" aria-label="Made by Ashton Dowling">Made by Ashton Dowling</div>
</body>
</html>
"""

    if not os.path.exists(html_path) or open(html_path, "r", encoding="utf-8").read() != html_code:
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html_code)

    # Clear folders on start
    for folder in (UPLOAD_DIR, OUTPUT_DIR):
        os.makedirs(folder, exist_ok=True)
        for filename in os.listdir(folder):
            file_path = os.path.join(folder, filename)
            if os.path.isfile(file_path):
                os.remove(file_path)
            elif os.path.isdir(file_path):
                shutil.rmtree(file_path)

    @app.route("/")
    def home():
        return render_template("upload.html", max_file_mb=MAX_FILE_MB)

    @app.route("/status")
    def status():
        job = request.args.get("job", "").strip()
        if not job:
            return {"error": "missing job"}, 400
        with PROGRESS_LOCK:
            state = PROGRESS.get(job)
        if not state:
            # Unknown job: return a neutral payload so the client doesn't break
            resp = {"status": "unknown", "pct": 0, "step": 0, "label": ""}
            return resp, 200
        return {
            "status": state.get("status", "running"),
            "pct": int(state.get("pct", 0)),
            "step": int(state.get("step", 0)),
            "label": state.get("label", "")
        }, 200

    @app.route("/cancel", methods=["POST"])
    def cancel_job():
        data = {}
        try:
            data = request.get_json(silent=True) or {}
        except Exception:
            data = {}
        job = (data.get("job") or request.form.get("job") or "").strip()
        if not job:
            return {"error": "missing job"}, 400
        with CANCELED_LOCK:
            CANCELED_JOBS.add(job)
        # Do NOT mark PROGRESS as 'done' here; the client handles UI reset on cancel.
        return {"ok": True}, 200
    @app.route("/smoke/local")
    def smoke_local():
        tex = r"""
    \documentclass[12pt]{article}
    \usepackage[a4paper,margin=20mm]{geometry}
    \begin{document}
    \section*{Smoke Test}
    If you can read this, Tectonic works. \(\int_0^1 x^2\,dx = \tfrac{1}{3}\).
    \end{document}
    """
        try:
            pdf = compile_tex_with_tectonic(tex)
            path = os.path.join(OUTPUT_DIR, "smoke.pdf")
            with open(path, "wb") as f:
                f.write(pdf)
            return send_from_directory(OUTPUT_DIR, "smoke.pdf", as_attachment=True)
        except Exception as e:
            return (f"LaTeX compile failed: {str(e)[:4000]}", 500)
    @app.route("/download/<kind>")
    def download(kind: str):
        filename = "questions.pdf" if kind == "questions" else ("answers.pdf" if kind == "answers" else None)
        if not filename:
            return "Invalid file requested", 400

        meta = _read_run_meta() or {}
        mode = meta.get("mode")
        available = set(meta.get("available", []))

        path = os.path.join(OUTPUT_DIR, filename)
        if not os.path.exists(path):
            if available and kind not in available:
                return (
                    f"The file for <em>{kind}.pdf</em> wasn’t produced in the last run. "
                    f"Available: {', '.join(sorted(available)) or 'none'}.",
                    404
                )
            # default fallback
            return "File not found. Generate it first.", 404

        # Normal send
        return send_from_directory(OUTPUT_DIR, filename, as_attachment=True)

    @app.route("/download/manifest")
    def download_manifest():
        meta = _read_run_meta() or {}
        # Don’t include anything sensitive; just mode/title/available/count/timestamp
        safe = {
            "mode": meta.get("mode"),
            "title": meta.get("title"),
            "available": meta.get("available", []),
            "n_items": meta.get("n_items"),
            "timestamp": meta.get("timestamp"),
        }
        return safe, 200
    def continue_mark_scheme(prev_text: str, start_idx: int, end_idx: int, questions_text: str, max_tokens: int) -> str:
        instr = (f"Continue the mark scheme starting from item {start_idx} to {end_idx}. "
                 "Use terse bullet points per item. Do not repeat previous items.\n\n"
                 f"Question paper:\n{questions_text}\n\n"
                 f"Previous mark scheme items:\n{prev_text}\n")
        return get_response(instr, "", model=main_model, max_tokens=max_tokens)

    def parse_tagged_bullets(text: str) -> list[str]:
        lines = [ln.strip() for ln in (text or "").splitlines()]
        tagged = [ln for ln in lines if ln.startswith("- [F")]
        if tagged:
            return tagged
        return [ln for ln in lines if ln.startswith("- ")]

    def interleave_round_robin(bullet_lists: list[list[str]], start_offset: int | None = None) -> list[str]:
        n = len(bullet_lists)
        if n == 0:
            return []
        order = list(range(n))
        if start_offset is None:
            start_offset = random.randrange(n)
        order = order[start_offset:] + order[:start_offset]
        positions = [0] * n
        total = sum(len(b) for b in bullet_lists)
        out = []
        while len(out) < total:
            for i in order:
                pos = positions[i]
                if pos < len(bullet_lists[i]):
                    out.append(bullet_lists[i][pos])
                    positions[i] += 1
        return out

    @app.route("/upload", methods=["POST"])
    def upload_file():
        # --- Stage 9: ensure we have a job id immediately ---
        job = (request.form.get("job") or str(uuid4())).strip()
        try:
            set_progress(job, 2, step=1, label="Starting upload")
            # --- Stage 13: auth + rate limit for /upload (progress-aware) ---
            ip = _client_ip()

            if BASIC_AUTH_ENABLED and not _auth_ok_for_request():
                # Finish the spinner politely using fail_progress
                return fail_progress(job, pct=96, step=1, label="Auth required",
                                     http_status=401, msg=ERR["auth_required"])

            if not _rate_allow("upload", ip, RATE_UPLOADS_PER_MIN):
                return fail_progress(job, pct=96, step=1, label="Rate limited",
                                     http_status=429, msg=ERR["rl_upload"])
            if "file[]" not in request.files:
                return fail_progress(job, pct=96, step=1, label="No files received",
                                     http_status=400, msg=ERR["no_file_part"])
            files = request.files.getlist("file[]")
            # --- Stage 2: new request shape & validation ---
            mode = (request.form.get("mode") or "exam").strip().lower()
            if mode not in VALID_MODES:
                return fail_progress(job, pct=96, step=1, label="Invalid mode",
                                     http_status=400, msg=ERR["invalid_mode"])

            title = (request.form.get("title") or "").strip()
            if not title:
                return fail_progress(job, pct=96, step=1, label="Missing title",
                                     http_status=400, msg=ERR["missing_title"])
            if len(title) > 80:
                return fail_progress(job, pct=96, step=1, label="Title too long",
                                     http_status=400, msg=ERR["long_title"])

            # difficulty may be "", which means "per-question overrides" in later stages.
            difficulty_raw = (request.form.get("difficulty") or "").strip().lower()
            difficulty_norm = difficulty_raw if difficulty_raw in DIFF_ALLOWED else None
            cfg = get_difficulty_profile(difficulty_norm or "medium")  # keep system stable for now

            # qcount is required only in exam mode (frontend enforces 1..30)
            qcount = parse_int(request.form.get("qcount"), default=0)
            if mode == "exam" and qcount <= 0:
                return fail_progress(job, pct=96, step=1, label="Missing question count",
                                     http_status=400, msg=ERR["missing_qcount"])
            if qcount > 30:
                return fail_progress(job, pct=96, step=1, label="Too many questions",
                                     http_status=400, msg=ERR["too_many_questions"])

            # Legacy per-type counts (still accepted for backward compatibility)
            num_long = parse_int(request.form.get("long"), 0)
            num_short = parse_int(request.form.get("short"), 0)
            num_mcq = parse_int(request.form.get("mcq"), 0)
            num_math = parse_int(request.form.get("math"), 0)

            # NOTE (Stage 2): We intentionally DO NOT require legacy counts anymore.
            # Generation still uses these values for now; Stage 3 will map qcount -> blueprint.
            # If all counts are 0, later stages will address generation logic.
            # --- Stage 3: construct blueprint ---
            legacy_total = num_long + num_short + num_mcq + num_math

            if legacy_total > 0:
                # Backward compatibility path (old clients still sending per-type counts)
                blueprint = blueprint_from_legacy_counts(num_long, num_short, num_mcq, num_math)
            else:
                # New contract: drive everything from qcount
                blueprint = build_default_blueprint(qcount)

            # Derive legacy counts from blueprint so the rest of the pipeline still works unchanged
            _bp_counts = counts_from_blueprint(blueprint)
            num_long = _bp_counts["Long"]
            num_short = _bp_counts["Short"]
            num_mcq = _bp_counts["MCQ"]
            num_math = _bp_counts["Math"]
            # --- Stage 4: read optional per-question overrides (exam mode only) ---
            if mode == "exam":
                raw_types = request.form.get("q_types")
                raw_topics = request.form.get("q_topics")
                raw_diffs = request.form.get("q_diffs")

                seq_types = _parse_seq_field(raw_types)
                seq_topics = _parse_seq_field(raw_topics)
                seq_diffs = _parse_seq_field(raw_diffs)

                N = len(blueprint)

                # If any override is provided, validate its length exactly equals qcount/blueprint length.
                def _len_ok(name, seq):
                    if seq is None:
                        return True
                    return len(seq) == N

                if not _len_ok("q_types", seq_types):
                    return fail_progress(job, pct=96, step=2, label="Bad overrides",
                                         http_status=422, msg=ERR["seq_len"].format(name="q_types", N=N))
                if not _len_ok("q_topics", seq_topics):
                    return fail_progress(job, pct=96, step=2, label="Bad overrides",
                                         http_status=422, msg=ERR["seq_len"].format(name="q_topics", N=N))
                if not _len_ok("q_diffs", seq_diffs):
                    return fail_progress(job, pct=96, step=2, label="Bad overrides",
                                         http_status=422, msg=ERR["seq_len"].format(name="q_diffs", N=N))

                # Validate contents; collect normalized forms (None if field absent)
                norm_types = None
                if seq_types is not None:
                    norm_types = []
                    for i, v in enumerate(seq_types, start=1):
                        t = _norm_type(v)
                        if not t:
                            return fail_progress(job, pct=96, step=2, label="Bad type override",
                                                 http_status=422, msg=ERR["invalid_qtype_at"].format(i=i, v=v))
                        norm_types.append(t)

                norm_diffs = None
                if seq_diffs is not None:
                    norm_diffs = []
                    for i, v in enumerate(seq_diffs, start=1):
                        d = _norm_diff(v)
                        if not d:
                            return fail_progress(job, pct=96, step=2, label="Bad difficulty override",
                                                 http_status=422, msg=ERR["invalid_diff_at"].format(i=i, v=v))
                        norm_diffs.append(d)

                clean_topics = None
                if seq_topics is not None:
                    clean_topics = []
                    for i, v in enumerate(seq_topics, start=1):
                        topic = str(v).strip()
                        if len(topic) > MAX_TOPIC_LEN:
                            return fail_progress(job, pct=96, step=2, label="Additional instructions too long",
                                                 http_status=422,
                                                 msg=ERR["topic_too_long_at"].format(i=i, max=MAX_TOPIC_LEN))
                        clean_topics.append(topic)

                # Apply overrides in-place
                if any(x is not None for x in (norm_types, clean_topics, norm_diffs)):
                    for i in range(N):
                        if norm_types is not None:
                            blueprint[i]["type"] = norm_types[i]
                        if clean_topics is not None and clean_topics[i]:
                            blueprint[i]["topic"] = clean_topics[i]
                        if norm_diffs is not None:
                            blueprint[i]["difficulty"] = norm_diffs[i]

                    # Re-derive legacy counts so the rest of the (pre-Stage-6) pipeline works unchanged
                    _bp_counts = counts_from_blueprint(blueprint)

                    set_progress(job, 34, step=2, label="Applying per-question overrides")
                    if is_canceled(job): return ("Canceled", 499)
            # Stage 11: enforce counts and totals up-front
            if len(files) > MAX_FILES:
                return fail_progress(job, pct=96, step=1, label="Too many files",
                                     http_status=400, msg=ERR["too_many_files"])

            total_bytes = 0
            seen_hashes = set()
            filepaths = []

            for f in files:
                if f.filename == "":
                    continue

                if not allowed_file(f.filename):
                    return fail_progress(job, pct=96, step=1, label="Invalid file format",
                                         http_status=415, msg=ERR["invalid_ext"].format(name=f.filename))

                # Stream to a temp file, compute sha256 and size
                hasher = hashlib.sha256()
                tmp = tempfile.NamedTemporaryFile(dir=UPLOAD_DIR, delete=False)
                size = 0
                try:
                    while True:
                        chunk = f.stream.read(64 * 1024)
                        if not chunk:
                            break
                        hasher.update(chunk)
                        tmp.write(chunk)
                        size += len(chunk)
                        if size > MAX_FILE_MB * 1024 * 1024:
                            tmp.close()
                            os.unlink(tmp.name)
                            return fail_progress(job, pct=96, step=1, label="File too large",
                                                 http_status=413, msg=ERR["file_too_big"])
                finally:
                    tmp.flush();
                    tmp.close()

                total_bytes += size
                if total_bytes > TOTAL_UPLOAD_MB * 1024 * 1024:
                    os.unlink(tmp.name)
                    return fail_progress(job, pct=96, step=1, label="Upload too large",
                                         http_status=413, msg=ERR["total_upload_too_big"])

                file_hash = hasher.hexdigest()
                if file_hash in seen_hashes:
                    # duplicate content; drop temp
                    os.unlink(tmp.name)
                    continue
                seen_hashes.add(file_hash)

                # Light content sniff (first 8 bytes)
                with open(tmp.name, "rb") as tfr:
                    head = tfr.read(8)

                ext = os.path.splitext(f.filename)[1].lower()
                ok = True
                if ext == ".pdf":
                    ok = _looks_pdf(head)
                    if ok:
                        # reject encrypted PDFs early
                        try:
                            with fitz.open(tmp.name) as d:
                                if getattr(d, "is_encrypted", False) and getattr(d, "needs_pass", False):
                                    os.unlink(tmp.name)
                                    return fail_progress(job, pct=96, step=2, label="Encrypted PDF",
                                                         http_status=422, msg=ERR["pdf_encrypted"])
                        except Exception:
                            pass  # will re-attempt in preprocessing or simply extract text fallback
                elif ext == ".rtf":
                    ok = _looks_rtf(head)
                elif ext in (".docx", ".pptx"):
                    ok = _looks_zip(head)
                    if ok:
                        kind = _office_zip_kind(tmp.name)
                        if (ext == ".docx" and kind != "docx") or (ext == ".pptx" and kind != "pptx"):
                            ok = False
                        if ok and not _zip_safety_ok(tmp.name):
                            os.unlink(tmp.name)
                            return fail_progress(job, pct=96, step=2, label="Unsafe Office archive",
                                                 http_status=422, msg=ERR["zip_bomb"])
                elif ext == ".txt":
                    # anything decodable as text is fine; we'll cap during extraction
                    ok = True
                else:
                    ok = False

                if not ok:
                    os.unlink(tmp.name)
                    return fail_progress(job, pct=96, step=1, label="Content/extension mismatch",
                                         http_status=415, msg=ERR["mime_mismatch"])

                # Move temp to final path
                # Move temp to final path (serialized to avoid Windows races)
                safe_name = secure_filename(f.filename)
                filepath = os.path.join(UPLOAD_DIR, safe_name)

                with UPLOAD_FS_LOCK:
                    if os.path.exists(filepath):
                        base, extn = os.path.splitext(safe_name)
                        safe_name = f"{base}-{file_hash[:8]}{extn}"
                        filepath = os.path.join(UPLOAD_DIR, safe_name)

                    # On Windows, AV/indexers may briefly lock the destination; retry a few times
                    moved = False
                    for _ in range(6):  # ~300ms total
                        try:
                            os.replace(tmp.name, filepath)
                            moved = True
                            break
                        except PermissionError:
                            time.sleep(0.05)

                    if not moved:
                        # Final fallback: unique suffix to avoid overwriting an in-use path
                        base, extn = os.path.splitext(safe_name)
                        unique = f"{base}-{uuid4().hex[:6]}{extn}"
                        filepath = os.path.join(UPLOAD_DIR, unique)
                        os.replace(tmp.name, filepath)

                filepaths.append(filepath)
            set_progress(job, 10, step=1, label="Processing files")
            if is_canceled(job): return ("Canceled", 499)
            if not filepaths:
                return fail_progress(job, pct=96, step=1, label="No valid files",
                                     http_status=400, msg=ERR["no_valid_files"])
            def _proc(i, p):
                return i, (preprocessing(p) or "")

            docs_idx_and_text = parallel_map(_proc, filepaths, max_workers=2)
            docs_idx_and_text.sort(key=lambda t: t[0])
            docs = [t[1] for t in docs_idx_and_text]
            set_progress(job, 25, step=2, label="Extracting text")
            if is_canceled(job): return ("Canceled", 499)
            # Stage 11: cap total extracted text across all docs
            running = 0
            capped_docs = []
            for t in docs:
                if running >= TOTAL_TEXT_CHAR_CAP:
                    capped_docs.append("")  # discard extras
                    continue
                room = TOTAL_TEXT_CHAR_CAP - running
                if len(t) <= room:
                    capped_docs.append(t)
                    running += len(t)
                else:
                    capped_docs.append(t[:room])
                    running += room
            docs = capped_docs
            if not any(docs):
                return fail_progress(job, pct=96, step=2, label="Extraction failed",
                                     http_status=422, msg=ERR["extraction_failed"])
            n_files = len(docs)
            if n_files == 0:
                return "No valid files to process", 400

            per_file_tokens = [fast_token_estimate(d) for d in docs]

            # --- Stage 8: qpaper mode (answers-only) ---
            raw_total_tokens = sum(per_file_tokens)
            raw_avg_tokens = max(1, raw_total_tokens // n_files)

            q_total = len(blueprint)

            T_SLA = 50.0
            HEADROOM = 5.0

            # Stage 5: derive caps strictly from the blueprint (works for qcount-only runs)
            n_out_q_cap, n_out_a_cap = estimate_output_token_caps_from_bp(blueprint)

            # Adaptive summary budget
            T_rem_for_sum, K_parallel, S_use = plan_summarization_sla(
                {}, n_files, raw_avg_tokens, n_out_q_cap, n_out_a_cap
            )
            set_progress(job, 35, step=2, label="Planning summaries")
            if is_canceled(job): return ("Canceled", 499)
            # Use SLA-driven S_use to decide summarize vs. full text
            R_thresh_exact = break_even_raw_avg_tokens_per_file(n_files, K_parallel, S_use)
            R_thresh_thumb = break_even_raw_avg_tokens_per_file_simple(S_use)
            # Use 0.8x threshold to favor summarization more often
            should_summarize = (raw_avg_tokens >= min(R_thresh_exact, R_thresh_thumb) * 0.8)

            N_in_Q_budget = max_input_tokens_for_main_questions(
                n_files=n_files,
                total_raw_tokens=raw_total_tokens,
                total_questions=q_total,
                n_out_q_tokens=n_out_q_cap,
                n_out_a_tokens=n_out_a_cap,
                n_summary_calls_parallel=K_parallel,
                raw_avg_tokens_per_file=raw_avg_tokens,
                S_summary_tokens_per_file=S_use if should_summarize else 0,
                T_total_sla=T_SLA,
                headroom_s=HEADROOM,
                T_sum_budget_s=T_rem_for_sum if should_summarize else 0.0,
            )
            N_in_Q_budget = max(ALWAYS_SAFE_MAIN_Q_INPUT_CAP, N_in_Q_budget)

            indices = list(range(n_files))
            random.shuffle(indices)
            docs_shuffled = [docs[i] for i in indices]

            if (not should_summarize) and (raw_total_tokens <= N_in_Q_budget):
                set_progress(job, 48, step=2, label="Using full text (within model budget)")
                if is_canceled(job): return ("Canceled", 499)
                text_for_main = "\n\n-----\n\n".join(docs_shuffled)
            else:
                set_progress(job, 48, step=2, label="Creating summaries (token-efficient)")
                if is_canceled(job): return ("Canceled", 499)
                TOK_PER_BULLET = 28
                # Scale bullets/tokens by difficulty (shorter for easy, longer for hard)
                target_bullets = max(8, min(36, int((S_use // TOK_PER_BULLET) * cfg["sum_token_scale"])))
                per_file_target = int(S_use * cfg["sum_token_scale"])

                prompt_tpl = (
                    "Extract exactly {target_bullets} key exam points from this material.\n"
                    "Format: '- [F{tag}] [fact/formula/concept]' ({words_each} words each)\n"
                    "Focus: definitions, formulas, processes, examples, and testable content.\n"
                    "{style_line}\n"
                    "Target output: ~{target_tokens} tokens total\n"
                    "Material:\n"
                )

                summaries = [None] * n_files
                with ThreadPoolExecutor(max_workers=K_parallel) as ex:
                    futs = []
                    for local_idx, doc in enumerate(docs_shuffled):
                        original_idx = indices[local_idx]
                        tag = original_idx + 1

                        adjusted_target = adaptive_summary_length(
                            fast_token_estimate(doc),
                            per_file_target
                        )

                        prompt = prompt_tpl.format(
                            tag=tag,
                            target_bullets=target_bullets,
                            words_each=cfg["sum_words_each"],
                            style_line=cfg["sum_style"],
                            target_tokens=adjusted_target,
                        )

                        futs.append((
                            local_idx,
                            ex.submit(get_response, prompt, doc, summary_model, adjusted_target, cfg["sum_temp"])
                        ))

                    for local_idx, fut in futs:
                        try:
                            s = fut.result()
                        except Exception:
                            s = ""
                        summaries[local_idx] = s or ""
                bullets_per_file = [parse_tagged_bullets(s) for s in summaries]
                for local_idx, blts in enumerate(bullets_per_file):
                    if not blts:
                        tag = indices[local_idx] + 1
                        lines = [ln.strip() for ln in (summaries[local_idx] or "").splitlines() if ln.strip()]
                        bullets_per_file[local_idx] = [f"- [F{tag}] {ln.lstrip('- ')}" for ln in lines[:target_bullets]]
                interleaved_bullets = interleave_round_robin(bullets_per_file, start_offset=None)
                text_for_main = "\n".join(interleaved_bullets)

            set_progress(job, 55, step=3, label="Summaries ready")
            if is_canceled(job): return ("Canceled", 499)
            # === Main questions call ===
            # === Main questions call (Stage 6: blueprint-driven) ===
            questions_needed = len(blueprint)

            # Build per-item prompt (use global difficulty only if provided; per-item diffs override)
            instruction_q = get_quality_question_instruction_from_blueprint(
                blueprint,
                global_difficulty=difficulty_norm  # None if user left difficulty blank
            )

            # Enhance material weighting for math if the blueprint is math-heavy
            bp_counts = counts_from_blueprint(blueprint)
            enhanced_material = enhance_math_content_for_questions(
                text_for_main,
                num_math=bp_counts["Math"],
                total_questions=questions_needed
            )

            max_tokens_q = n_out_q_cap
            questions = get_response(
                instruction_q,
                enhanced_material,
                model=main_model,
                max_tokens=max_tokens_q,
                temperature=cfg["q_temp"]
            )
            questions = questions.replace('*', '')
            questions_needed = len(blueprint)
            q_items = split_numbered_items(questions)
            if len(q_items) < questions_needed:
                start_missing = len(q_items) + 1
                end_missing = questions_needed
                cont = get_response(
                    continue_numbered_list_from_blueprint(
                        prev_text=questions,
                        start_idx=start_missing,
                        end_idx=end_missing,
                        material=text_for_main,
                        blueprint=blueprint,
                        global_difficulty=difficulty_norm
                    ),
                    "",
                    model=main_model,
                    max_tokens=int(n_out_q_cap * 0.5),
                    temperature=cfg["q_temp"]
                )
                q_items += split_numbered_items(cont)[: (questions_needed - len(q_items))]
            q_items = clamp_items(q_items, questions_needed)
            q_items = strip_headers_from_items(q_items)

            # LaTeX-safe transform + fixers
            def process_latex_item(item: str) -> str:
                # 0) Normalize text & basic mapping
                item = latex_backup_translate(normalize_for_latex(item))

                # 1) Bring any math-y fragments INTO math first
                item = _wrap_naked_math(item)

                # 2) Inside-math structural fixes (vectors/hats and bad frac forms)
                item = _fix_veclike_args_in_math(item)
                item = _fix_text_macros_in_math(item)
                item = _fix_frac_forms_in_math(item)

                # 3) Convert inline a/b only INSIDE math into \frac{a}{b}
                item = convert_slashes_only_inside_math(item)

                # 4) Your existing malformed \frac + backslash sanitizers
                item = _fix_malformed_frac_text(item)
                item = _sanitize_backslashes(item)

                return item

            q_items = parallel_map(lambda i, item: process_latex_item(item), q_items, max_workers=4)
            set_progress(job, 70, step=4, label="Generating questions")
            if is_canceled(job): return ("Canceled", 499)
            # --- ANSWERS ---
            # Stage 7: blueprint-aligned mark scheme generation
            answers_needed = len(blueprint)
            instruction_a = get_quality_answer_instruction_from_blueprint(blueprint)
            max_tokens_a = n_out_a_cap

            set_progress(job, 76, step=4, label="Generating mark scheme")
            if is_canceled(job): return ("Canceled", 499)
            answers = get_response(
                instruction_a,
                questions,  # pass the question paper as before
                model=main_model,
                max_tokens=max_tokens_a,
                temperature=0.2
            )
            answers = answers.replace('*', '')

            a_items = split_numbered_items(answers)

            # If under-filled, continue with a blueprint-aware prompt
            if len(a_items) < answers_needed:
                start_missing = len(a_items) + 1
                end_missing = answers_needed
                cont_a_instr = continue_mark_scheme_from_blueprint(
                    prev_text=answers,
                    start_idx=start_missing,
                    end_idx=end_missing,
                    questions_text=questions,
                    blueprint=blueprint,
                )
                cont_a = get_response(
                    cont_a_instr,
                    "",
                    model=main_model,
                    max_tokens=int(n_out_a_cap * 0.5),
                    temperature=0.2
                )
                a_items += split_numbered_items(cont_a)[: (answers_needed - len(a_items))]

            # Enforce 1:1 alignment with blueprint length
            if len(a_items) > answers_needed:
                a_items = a_items[:answers_needed]

            a_items = strip_headers_from_items(a_items)
            # After you build a_items:
            # --- ANSWERS SANITIZE (mirror questions + new sqrt/frac fixes) ---
            a_items = [latex_backup_translate(normalize_for_latex(it)) for it in a_items]
            a_items = [_wrap_naked_math(it) for it in a_items]
            a_items = [_fix_veclike_args_in_math(it) for it in a_items]
            a_items = [_fix_frac_forms_in_math(it) for it in a_items]
            a_items = [_fix_sqrt_args_in_math(it) for it in a_items]  # NEW
            a_items = [_fix_frac_sqrt_edgecases_in_math(it) for it in a_items]  # NEW
            a_items = [convert_slashes_only_inside_math(it) for it in a_items]
            a_items = [_fix_malformed_frac_text(it) for it in a_items]
            a_items = [_sanitize_backslashes(it) for it in a_items]
            set_progress(job, 85, step=5, label="Mark scheme ready")
            if is_canceled(job): return ("Canceled", 499)
            # Build & compile
            default_title = (title or "").strip() or "Mock Exam Paper"
            q_tex = tex_from_items(q_items, default_title.title())
            set_progress(job, 90, step=5, label="Compiling mark scheme")
            if is_canceled(job): return ("Canceled", 499)
            a_tex = tex_from_items(a_items, default_title.title() + " Answers")
            q_tex = _emergency_tex_sanitize(q_tex)
            a_tex = _emergency_tex_sanitize(a_tex)
            set_progress(job, 90, step=5, label="Compiling PDFs")
            if is_canceled(job): return ("Canceled", 499)
            q_path = os.path.join(OUTPUT_DIR, "questions.pdf")
            a_path = os.path.join(OUTPUT_DIR, "answers.pdf")

            # compile sequentially, directly to disk (no giant in-RAM PDFs)
            compile_or_repair_to_path(q_tex, q_path)
            compile_or_repair_to_path(a_tex, a_path)
            _write_run_meta(
                mode="exam",
                title=default_title.title(),
                available=["questions", "answers"],
                extra={"n_items": len(q_items)}
            )
            set_progress(job, 100, step=5, label="Ready to download", status="done")

            links = (
                '<p style="margin-top:10px;">'
                '<a href="/download/questions" target="_blank">Download questions.pdf</a> &middot; '
                '<a href="/download/answers" target="_blank">Download answers.pdf</a>'
                '</p>'
            )
            return links, 200
        except Exception as e:
            # Mark as done (so the UI can finish), then return 500
            try:
                app.logger.exception("Upload failed [job=%s]: %s", job, e)
            except Exception:
                pass
            return fail_progress(job, pct=98, step=5, label="Internal error",
                                 http_status=500, msg=ERR["internal"])
    # Stage 14: log effective config (no secrets)
    app.logger.info("Config: MODELS main=%s summary=%s | TECTONIC_TIMEOUT=%ss", main_model, summary_model, TECTONIC_TIMEOUT)
    app.logger.info("Limits: MAX_FILES=%d, MAX_FILE_MB=%d, TOTAL_UPLOAD_MB=%d", MAX_FILES, MAX_FILE_MB, TOTAL_UPLOAD_MB)
    app.logger.info("Caps: SUMMARY %d (%d..%d), Q_IN=%d, Q_OUT=%d, A_OUT=%d",
                    RECOMMENDED_SUMMARY_TOKENS, SUMMARY_TOKENS_HARD_MIN, SUMMARY_TOKENS_HARD_MAX,
                    ALWAYS_SAFE_MAIN_Q_INPUT_CAP, TARGET_MAX_N_OUT_Q, TARGET_MAX_N_OUT_A)

    # --- Stage 15: error handlers ---

    @app.errorhandler(413)
    def _too_large(e):
        # Body too large (Flask rejected before hitting /upload)
        return "Total upload exceeds server limit.", 413

    @app.errorhandler(429)
    def _too_many(e):
        return "Too many requests.", 429

    @app.errorhandler(401)
    def _auth_needed(e):
        # For non-/upload endpoints; /upload already uses fail_progress
        r = Response("Authentication required.", 401)
        r.headers["WWW-Authenticate"] = 'Basic realm="ExamApp"'
        return r

    @app.errorhandler(404)
    def _nf(e):
        return "Not found.", 404

    @app.errorhandler(500)
    def _ise(e):
        # Avoid leaking stack traces in prod; logs already capture exception
        return "Internal server error.", 500
    return app

if __name__ == "__main__":
    # Stage 1: env-driven dev run (prod should use a WSGI server like gunicorn)
    logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(name)s: %(message)s")
    debug = os.getenv("FLASK_DEBUG", "0").lower() in ("1","true","yes","on")
    host  = os.getenv("APP_HOST", "0.0.0.0")
    port  = int(os.getenv("APP_PORT", "5000"))
    app = website()  # website() should return `app` (see next change)
    app.run(debug=debug, host=host, port=port, threaded=True)

